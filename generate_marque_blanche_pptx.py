#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Génère une présentation PowerPoint HeliaPV — Offre Marque Blanche
Cible : Fabricants · Distributeurs · Développeurs de projets
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

OUTPUT = Path("HeliaPV_MarqueBlanche_2026.pptx")

# ── Palette ────────────────────────────────────────────────────────────────────
C_BG_DARK  = RGBColor(0x0A, 0x0E, 0x1A)
C_BG_CARD  = RGBColor(0x11, 0x18, 0x27)
C_BG_CARD2 = RGBColor(0x14, 0x1C, 0x2E)
C_ACCENT   = RGBColor(0xF5, 0x9E, 0x0B)   # doré solaire
C_ACCENT2  = RGBColor(0x63, 0x66, 0xF1)   # indigo
C_GREEN    = RGBColor(0x10, 0xB9, 0x81)   # vert succès
C_WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
C_MUTED    = RGBColor(0x94, 0xA3, 0xB8)
C_FAB      = RGBColor(0xEF, 0x44, 0x44)   # rouge – fabricants
C_DIST     = RGBColor(0x06, 0xB6, 0xD4)   # cyan – distributeurs
C_DEV      = RGBColor(0x8B, 0x5C, 0xF6)   # violet – développeurs

W = Inches(13.33)
H = Inches(7.5)

# ── Utilitaires ────────────────────────────────────────────────────────────────

def rgb(r, g, b):
    return RGBColor(r, g, b)


def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_text(slide, text, left, top, width, height,
             font_size=18, bold=False, color=C_WHITE,
             align=PP_ALIGN.LEFT, italic=False):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Segoe UI"
    return txb


def add_bullet_list(slide, bullets, left, top, width, height,
                    font_size=11, color=C_WHITE, accent=C_ACCENT):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for b in bullets:
        para = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        para.space_before = Pt(2)
        r1 = para.add_run()
        r1.text = "▸  "
        r1.font.color.rgb = accent
        r1.font.size = Pt(font_size)
        r1.font.name = "Segoe UI"
        r2 = para.add_run()
        r2.text = b
        r2.font.color.rgb = color
        r2.font.size = Pt(font_size)
        r2.font.name = "Segoe UI"
    return txb


def add_pill(slide, text, left, top, bg=C_ACCENT, fg=C_BG_DARK, w=Inches(3.5), h=Inches(0.3)):
    add_rect(slide, left, top, w, h, bg)
    add_text(slide, text, left + Inches(0.1), top, w - Inches(0.1), h,
             font_size=8.5, bold=True, color=fg)


def slide_n_label(slide, n):
    add_text(slide, str(n), Inches(12.9), Inches(7.1), Inches(0.4), Inches(0.3),
             font_size=10, color=C_MUTED, align=PP_ALIGN.RIGHT)


def header_bar(slide, title, category=None, category_color=C_ACCENT):
    add_rect(slide, 0, 0, Inches(0.08), H, C_ACCENT)
    add_rect(slide, 0, H - Inches(0.08), W, Inches(0.08), C_ACCENT2)
    if category:
        add_pill(slide, category, Inches(0.3), Inches(0.18),
                 bg=category_color, fg=C_BG_DARK)
    y_title = Inches(0.55) if category else Inches(0.3)
    add_text(slide, title, Inches(0.3), y_title, Inches(12.7), Inches(0.65),
             font_size=24, bold=True, color=C_WHITE)
    add_rect(slide, Inches(0.3), y_title + Inches(0.7), Inches(12.7), Inches(0.03), C_ACCENT)


# ── SLIDES DATA ────────────────────────────────────────────────────────────────

SLIDES = [

    # ── 1 · Couverture ─────────────────────────────────────────────────────────
    {
        "type": "cover",
        "title": "HeliaPV",
        "subtitle": "Votre plateforme solaire\nen marque blanche",
        "tagline": "Offrez à vos clients la meilleure solution SaaS du marché — sous votre propre marque",
        "version": "Offre Partenariat · Mars 2026",
        "targets": ["Fabricants de matériel PV", "Distributeurs & Grossistes", "Développeurs de projets"],
    },

    # ── 2 · Sommaire ───────────────────────────────────────────────────────────
    {
        "type": "toc",
        "title": "Au programme",
        "sections": [
            ("1", "L'opportunité marché",          "Croissance PV en France, besoin d'outils"),
            ("2", "Pour qui ?",                     "Fabricants · Distributeurs · Développeurs"),
            ("3", "Fabricants de matériel PV",      "Augmentez votre valeur ajoutée"),
            ("4", "Distributeurs & Grossistes",     "Fidélisez vos réseaux d'installateurs"),
            ("5", "Développeurs de projets",        "Déployez votre propre outil de prospection"),
            ("6", "La plateforme complète",         "Tous les modules inclus"),
            ("7", "Marque blanche totale",          "Logo · Couleurs · Domaine · URL"),
            ("8", "Intégration technique & API",    "Multi-tenant · SSO · Webhooks"),
            ("9", "Modèle économique",              "Licence OEM · Revenue share · Custom"),
            ("10", "Mise en service en 30 jours",   "Onboarding · Formation · Support"),
        ],
    },

    # ── 3 · Opportunité marché ─────────────────────────────────────────────────
    {
        "type": "stats",
        "category": "L'OPPORTUNITÉ MARCHÉ",
        "title": "Le marché solaire français : une croissance structurelle",
        "stats": [
            ("90 GWc", "Objectif France 2030\n(×4 vs 2023)"),
            ("3 000 €/kWc", "Prix moyen d'une\ninstallation résidentielle"),
            ("15 000+", "Installateurs&\ndistributeurs actifs"),
            ("2 M", "Contacts propriétaires\ndans notre base MAJIC"),
        ],
        "bullets": [
            "La loi Énergie-Climat (2023) impose l'accélération des EnR : les AO CRE s'intensifient",
            "Les fabricants et distributeurs cherchent à se différencier dans un marché mature et concurrentiel",
            "Les installateurs ont besoin d'outils de prospection performants — aujourd'hui ils n'en ont pas",
            "Aucune solution SaaS complète de bout-en-bout n'est disponible en marque blanche sur le marché français",
            "HeliaPV couvre l'intégralité du cycle : prospection → calepinage → docs → chantier → mise en service",
        ],
    },

    # ── 4 · Pour qui ───────────────────────────────────────────────────────────
    {
        "type": "persona",
        "title": "Une offre taillée pour trois profils partenaires",
        "personas": [
            {
                "icon": "🏭",
                "title": "Fabricants\nde matériel PV",
                "color": C_FAB,
                "points": [
                    "Différenciez votre catalogue",
                    "Générez des leads pour vos\nrevendeurs",
                    "Valorisez vos modules dans\nles simulations calepinage",
                    "Accédez aux données de\nconsommation marché",
                    "Renforcez la fidélité de vos\nrevendeurs / intégrateurs",
                ],
            },
            {
                "icon": "🏪",
                "title": "Distributeurs\n& Grossistes",
                "color": C_DIST,
                "points": [
                    "Offrez un outil métier premium\nà vos clients installateurs",
                    "Créez un écosystème digital\nautour de votre marque",
                    "Suivez l'activité commerciale\nde votre réseau",
                    "Proposez une formation\nintégrée à vos revendeurs",
                    "Générez un revenu récurrent\nen SaaS",
                ],
            },
            {
                "icon": "⚡",
                "title": "Développeurs\nde projets",
                "color": C_DEV,
                "points": [
                    "Outil de prospection exclusif\nsous votre marque",
                    "Cartographie GIS + LiDAR HD\npour vos équipes terrain",
                    "CRM adapté à votre\nprocessus de développement",
                    "AO CRE automatisés\n(bâtiment, agrivoltaïque)",
                    "Accès API pour vos outils\ninternes (BI, ERP, GED)",
                ],
            },
        ],
    },

    # ── 5 · Fabricants ─────────────────────────────────────────────────────────
    {
        "type": "plain",
        "category": "FABRICANTS DE MATÉRIEL PV",
        "cat_color": C_FAB,
        "title": "Valorisez vos produits dans chaque projet installé",
        "description": (
            "En marque blanche HeliaPV, vos modules, onduleurs ou systèmes de fixation "
            "sont visibles et utilisés par vos revendeurs à chaque nouveau projet."
        ),
        "bullets_left": [
            "Catalogue produits intégré : vos modules & onduleurs disponibles dans l'outil de calepinage",
            "Simulation PVGIS automatique avec vos références techniques (Pmax, Isc, Voc, température)",
            "Plans de string, schémas unifilaires et CERFA générés avec votre matériel",
            "Proposition commerciale PDF avec vos visuels, prix et références produit",
            "Reporting agrégé : MW dimensionnés, volumes, taux d'adoption par région",
        ],
        "bullets_right": [
            "Vos installateurs restent dans votre écosystème digital plutôt que chez un concurrent",
            "Dashboard fabricant : suivi de l'utilisation de vos produits en temps réel",
            "Exportation leads qualifiés : projets utilisant vos modules, prêts à commander",
            "Module de certifications intégré : documentation NF, CE, IEC facilement accessible",
            "Connexion API ERP / PIM : synchronisation automatique de votre catalogue produit",
        ],
        "cta": "Vos modules dans 100% des projets de vos revendeurs, sans friction.",
    },

    # ── 6 · Distributeurs ──────────────────────────────────────────────────────
    {
        "type": "plain",
        "category": "DISTRIBUTEURS & GROSSISTES",
        "cat_color": C_DIST,
        "title": "Transformez votre offre produit en écosystème digital",
        "description": (
            "Proposez à chaque installateur de votre réseau un outil SaaS complet "
            "sous votre marque. Un avantage concurrentiel décisif face aux pure players du digital."
        ),
        "bullets_left": [
            "Plateforme SaaS à votre nom : domaine personnalisé, logo, charte graphique",
            "Gestion multi-comptes : un admin par agence ou région, une vue consolidée pour vous",
            "Votre catalogue produit pré-chargé dans les devis automatiques de vos installateurs",
            "Génération automatique de devis reliés à vos prix grossiste : marge protégée",
            "Formation intégrée : tutoriels vidéo, base de connaissances sous votre marque",
        ],
        "bullets_right": [
            "Indicateurs réseau : qui utilise l'outil, combien de projets, volumes commandés",
            "Programme fidélité : points cumulés sur les projets traités dans la plateforme",
            "Contrat d'exclusivité territoriale possible selon votre zone de chalandise",
            "Revenu récurrent SaaS : vous facturez vos installateurs, nous vous facturons",
            "Support mutualisé : notre équipe technique répond en votre nom (marque blanche totale)",
        ],
        "cta": "Votre réseau, votre data, votre marque — une plateforme prête en 30 jours.",
    },

    # ── 7 · Développeurs ───────────────────────────────────────────────────────
    {
        "type": "plain",
        "category": "DÉVELOPPEURS DE PROJETS",
        "cat_color": C_DEV,
        "title": "Votre outil de développement exclusif, sous votre propre marque",
        "description": (
            "Que vous développiez des centrales au sol, des toitures industrielles ou des projets "
            "agrivoltaïques, HeliaPV en marque blanche devient votre avantage concurrentiel interne."
        ),
        "bullets_left": [
            "Cartographie GIS propriétaire : couches RPG, HTA, Friches, ZAER — scoring automatique",
            "Base propriétaires MAJIC 2M+ contacts : prospection foncière directe depuis la carte",
            "LiDAR HD IGN COPC : modélisation 3D toiture avec pentes, orientations, masques solaires",
            "Simulation P50/P90 PVGIS 8760h : dossiers bancaires et AO CRE directement exploitables",
            "Autoconsommation collective : identification des périmètres et participants en 1 clic",
        ],
        "bullets_right": [
            "CRM dédié développeurs : suivi des phases foncière, administrative, technique, financière",
            "AO bâtiment >500 kWc : dossier CRE4 constitué automatiquement avec pièces justificatives",
            "GED versionn\u00e9e : tous les documents du projet horodatés et stockés dans la plateforme",
            "API REST complète : intégration dans votre ERP, BI (Power BI, Tableau), GIS interne",
            "Multi-projet : suivez des centaines de projets simultanément avec une seule interface",
        ],
        "cta": "De la détection de site à la mise en service — sans changer d'outil.",
    },

    # ── 8 · Plateforme complète ────────────────────────────────────────────────
    {
        "type": "features_grid",
        "category": "LA PLATEFORME COMPLÈTE",
        "title": "16 modules inclus dans la licence marque blanche",
        "modules": [
            ("🗺️",  "Cartographie GIS",          "Leaflet + IGN + GeoServer"),
            ("🔍",  "Prospection territoriale",   "RPG · Friches · Postes HTA"),
            ("📊",  "CRM intégré",                "Pipeline · KPI · Relances"),
            ("📐",  "Calepinage 3D",              "PVGIS · Ombrage · Plan masse"),
            ("⚡",  "Simulation autoconsommation","8760h · TRI · VAN · 25 ans"),
            ("📡",  "LiDAR HD IGN",              "COPC · MNS · Vue 3D WebGL"),
            ("📄",  "Documents auto",             "CERFA · Schéma · Unifilaire"),
            ("💼",  "Propositions commerciales",  "PDF · Pricing · Financement"),
            ("🏗️",  "Suivi de chantier",         "IEC 62446 · DOE · CONSUEL"),
            ("🤖",  "IA Helia",                  "LLM solaire · Groq · LLama"),
            ("📋",  "AO PV bâtiment",            "CRE4 · >500 kWc · CERFA"),
            ("🌐",  "Autoconso collective",       "ACC · Participants · Clés"),
            ("📈",  "Rapports commune/GPS",       "PLU · GéoRisques · Enedis"),
            ("🔐",  "Auth multi-rôles",          "Admin · Premium · Trial"),
            ("⚙️",  "Backoffice admin",          "Licences · Quotas · Logs"),
            ("🔌",  "API REST complète",         "Webhooks · SSE · JSON"),
        ],
    },

    # ── 9 · Marque blanche totale ──────────────────────────────────────────────
    {
        "type": "plain",
        "category": "PERSONNALISATION MARQUE BLANCHE",
        "cat_color": C_GREEN,
        "title": "100% à votre image — vos clients ne voient jamais HeliaPV",
        "description": (
            "La plateforme est intégralement personnalisable. Votre identité visuelle, "
            "votre domaine, vos couleurs, vos emails, vos documents — aucune référence à HeliaPV."
        ),
        "bullets_left": [
            "Domaine personnalisé : app.votreentreprise.fr ou solar.votregroupe.com",
            "Logo en haute résolution sur toutes les interfaces, emails, PDF et exports",
            "Palette couleurs ajustée à votre charte graphique (primaire, secondaire, accent)",
            "Favicon, titre de page, meta descriptions SEO à votre nom",
            "Emails transactionnels : confirmation, relance, rapport — expéditeur @votredomaine.fr",
        ],
        "bullets_right": [
            "Propositions commerciales PDF avec entête à votre logo, adresse, CGV",
            "CERFA et documents techniques : cartouche à votre nom et coordonnées",
            "Page de connexion personnalisée : fond, slogan, visuels au choix",
            "Mentions légales, CGU et politique de confidentialité à votre nom",
            "Suppression totale de toute référence à HeliaPV / AgriWeb dans l'interface",
        ],
        "cta": "Vos clients voient votre marque — nous restons invisibles.",
    },

    # ── 10 · Intégration technique ─────────────────────────────────────────────
    {
        "type": "tech",
        "title": "Architecture marque blanche — Intégration technique",
        "layers": [
            ("Multi-tenant", [
                "Isolation totale des données par tenant",
                "Sous-domaine dédié par partenaire",
                "Configuration par tenant en BDD",
                "Quotas & limites personnalisés",
                "Logs séparés par organisation",
                "Backup isolé par partenaire",
            ]),
            ("API & Intégrations", [
                "API REST JSON documentée (Swagger)",
                "Webhooks configurables (projet, doc, CRM)",
                "SSO SAML 2.0 / OAuth2 (Azure AD, etc.)",
                "Export CSV / Excel / GeoJSON / KML",
                "Intégration ERP via endpoints dédiés",
                "Streaming SSE pour prospections longues",
            ]),
            ("Déploiement", [
                "SaaS mutualisé (Railway/Docker) — livraison immédiate",
                "Instance dédiée possible (VPS, cloud privé)",
                "On-premise sur vos serveurs si requis",
                "CI/CD GitHub Actions — mises à jour transparentes",
                "Certifications ISO 27001 en cours (audit 2026)",
                "RGPD : données hébergées en France (OVH/Scaleway)",
            ]),
            ("Support dédié", [
                "Référent technique dédié par partenaire",
                "SLA 99,5% de disponibilité garanti",
                "Bac à sable (sandbox) pré-intégration",
                "Documentation API & SDK Python/JS",
                "Formations onboarding (visio + présentiel)",
                "Hotline prioritaire H+4 en production",
            ]),
        ],
    },

    # ── 11 · Modèle économique ─────────────────────────────────────────────────
    {
        "type": "pricing",
        "category": "MODÈLE ÉCONOMIQUE",
        "title": "Trois formules de partenariat adaptées à votre cas d'usage",
        "plans": [
            {
                "name": "OEM Lite",
                "color": C_DIST,
                "price": "Sur devis",
                "subtitle": "À partir de 5 licences",
                "ideal": "Distributeurs régionaux\nFabricants spécialisés",
                "includes": [
                    "Marque blanche complète",
                    "Domaine personnalisé",
                    "10 modules inclus",
                    "Support Standard (J+2)",
                    "Mises à jour incluses",
                    "1 formation onboarding",
                ],
                "highlight": False,
            },
            {
                "name": "OEM Pro",
                "color": C_ACCENT,
                "price": "Sur devis",
                "subtitle": "À partir de 20 licences",
                "ideal": "Distributeurs nationaux\nFabricants tier-1",
                "includes": [
                    "Tout OEM Lite +",
                    "16 modules complets",
                    "API REST + Webhooks",
                    "Support Prioritaire (H+4)",
                    "Sandbox dédié",
                    "Dashboard partenaire",
                    "Revenue share possible",
                    "2 formations + docs custom",
                ],
                "highlight": True,
            },
            {
                "name": "OEM Enterprise",
                "color": C_DEV,
                "price": "Sur devis",
                "subtitle": "Déploiement sur mesure",
                "ideal": "Grands comptes\nDéveloppeurs industriels",
                "includes": [
                    "Tout OEM Pro +",
                    "Instance dédiée / on-premise",
                    "SSO SAML / Azure AD",
                    "Développements spécifiques",
                    "SLA 99,9% contractuel",
                    "Référent dédié permanent",
                    "Formation équipes terrain",
                    "Portage IP négociable",
                ],
                "highlight": False,
            },
        ],
        "note": "Pas de frais d'entrée sur OEM Lite et Pro · Contractualisation en 5 jours · Mise en prod en 30 jours",
    },

    # ── 12 · Mise en service en 30 jours ──────────────────────────────────────
    {
        "type": "timeline",
        "category": "ONBOARDING PARTENAIRE",
        "title": "De la signature au lancement : 30 jours",
        "steps": [
            ("Sem. 1", "Contractualisation & Kick-off",
             "Signature NDA + contrat OEM\nKick-off call avec référent dédié\nDéfinition périmètre & calendrier"),
            ("Sem. 2", "Personnalisation & Config",
             "Intégration logo, charte, domaine\nChargement catalogue produits\nCréation comptes admin & test"),
            ("Sem. 3", "Intégration & Tests",
             "Connexion API / SSO si requis\nTests fonctionnels en sandbox\nValidation conjointe de la plateforme"),
            ("Sem. 4", "Formation & Go-live",
             "Formation équipes (visio 2h)\nMise en production\nSupport renforcé J+30 post-launch"),
        ],
        "after": "Après le lancement : mises à jour automatiques, nouveaux modules inclus, roadmap partagée trimestrielle",
    },

    # ── 13 · Conclusion & CTA ──────────────────────────────────────────────────
    {
        "type": "conclusion",
        "title": "HeliaPV Marque Blanche — Passez à l'action",
        "points": [
            "🏷️  Plateforme complète sous votre marque — aucune référence HeliaPV visible",
            "🔌  Intégration API · SSO · Multi-tenant · Webhooks en standard",
            "⚡  Mise en production en 30 jours — onboarding clé en main",
            "📦  16 modules inclus : prospection → calepinage → docs → IA → chantier",
            "💶  Modèle OEM Lite / Pro / Enterprise — adapté à votre volume",
            "📡  LiDAR HD IGN + PVGIS 8760h — technologie introuvable ailleurs",
            "🤖  IA Helia (LLM solaire) — différenciant immédiat pour vos utilisateurs",
            "🇫🇷  Données hébergées en France · RGPD · Certif. ISO 27001 en cours",
        ],
        "url": "https://app.heliapv.fr",
        "contact": "yann.laurent@heliapv.fr · 06 21 16 55 85",
    },
]


# ── GÉNÉRATEURS DE SLIDES ──────────────────────────────────────────────────────

def make_cover(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, C_BG_DARK)

    # Bande accent gauche
    add_rect(slide, 0, 0, Inches(0.1), H, C_ACCENT)

    # Forme soleil décorative (coin haut droit)
    sun = slide.shapes.add_shape(1, Inches(9.8), Inches(-0.3), Inches(4), Inches(4))
    sun.fill.solid()
    sun.fill.fore_color.rgb = rgb(0xF5, 0x9E, 0x0B)
    sun.line.fill.background()
    sun2 = slide.shapes.add_shape(1, Inches(9.95), Inches(-0.15), Inches(3.7), Inches(3.7))
    sun2.fill.solid()
    sun2.fill.fore_color.rgb = C_BG_DARK
    sun2.line.fill.background()

    # Logo
    add_text(slide, "☀️  HeliaPV", Inches(0.5), Inches(1.0), Inches(6), Inches(0.7),
             font_size=13, bold=True, color=C_ACCENT)

    # Titre
    txb = slide.shapes.add_textbox(Inches(0.5), Inches(1.7), Inches(9), Inches(2.0))
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = data["subtitle"]
    r.font.size = Pt(44)
    r.font.bold = True
    r.font.color.rgb = C_WHITE
    r.font.name = "Segoe UI"

    # Tagline
    add_text(slide, data["tagline"], Inches(0.5), Inches(3.85), Inches(9.5), Inches(0.55),
             font_size=13, color=C_MUTED, italic=True)

    # Cibles
    TARGET_COLORS = [C_FAB, C_DIST, C_DEV]
    for i, t in enumerate(data["targets"]):
        x = Inches(0.5 + i * 4.1)
        add_rect(slide, x, Inches(4.6), Inches(3.8), Inches(0.38), TARGET_COLORS[i])
        add_text(slide, t, x + Inches(0.12), Inches(4.6), Inches(3.65), Inches(0.38),
                 font_size=11, bold=True, color=C_BG_DARK)

    # Version
    add_text(slide, data["version"], Inches(0.5), Inches(6.85), Inches(5), Inches(0.3),
             font_size=9, color=C_MUTED)

    # Bande basse
    add_rect(slide, 0, H - Inches(0.08), W, Inches(0.08), C_ACCENT2)
    add_text(slide, "app.heliapv.fr", Inches(9.5), Inches(6.85), Inches(3.6), Inches(0.3),
             font_size=10, color=C_ACCENT, align=PP_ALIGN.RIGHT)
    print("  ✅  Slide 1 : Couverture")


def make_toc(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, C_BG_DARK)
    add_rect(slide, 0, 0, Inches(0.08), H, C_ACCENT)
    add_text(slide, data["title"], Inches(0.3), Inches(0.25), Inches(12), Inches(0.6),
             font_size=26, bold=True, color=C_WHITE)
    add_rect(slide, Inches(0.3), Inches(0.9), Inches(12.7), Inches(0.04), C_ACCENT)

    cols = [data["sections"][:5], data["sections"][5:]]
    for ci, col in enumerate(cols):
        x = Inches(0.4 + ci * 6.4)
        for i, (num, title, sub) in enumerate(col):
            y = Inches(1.05 + i * 1.15)
            add_rect(slide, x, y, Inches(0.42), Inches(0.42), C_ACCENT)
            add_text(slide, num, x, y + Inches(0.05), Inches(0.42), Inches(0.32),
                     font_size=13, bold=True, color=C_BG_DARK, align=PP_ALIGN.CENTER)
            add_text(slide, title, x + Inches(0.52), y, Inches(5.6), Inches(0.3),
                     font_size=11.5, bold=True, color=C_WHITE)
            add_text(slide, sub, x + Inches(0.52), y + Inches(0.3), Inches(5.6), Inches(0.25),
                     font_size=8.5, color=C_MUTED, italic=True)
    print("  ✅  Slide 2 : Sommaire")


def make_stats(prs, data, n):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, C_BG_DARK)
    header_bar(slide, data["title"], data.get("category"), C_ACCENT)

    # 4 KPI boxes
    stats = data["stats"]
    box_w = Inches(3.0)
    box_h = Inches(1.7)
    for i, (val, lbl) in enumerate(stats):
        x = Inches(0.35 + i * 3.25)
        y = Inches(1.4)
        add_rect(slide, x, y, box_w, box_h, C_BG_CARD2)
        add_text(slide, val, x, y + Inches(0.1), box_w, Inches(0.75),
                 font_size=28, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
        add_text(slide, lbl, x, y + Inches(0.85), box_w, Inches(0.75),
                 font_size=9.5, color=C_MUTED, align=PP_ALIGN.CENTER)

    # Bullets insights
    add_text(slide, "Pourquoi maintenant ?", Inches(0.35), Inches(3.25),
             Inches(12.7), Inches(0.3), font_size=12, bold=True, color=C_WHITE)
    add_bullet_list(slide, data["bullets"], Inches(0.35), Inches(3.6),
                    Inches(12.6), Inches(3.5), font_size=11, accent=C_ACCENT)
    slide_n_label(slide, n)
    print(f"  ✅  Slide {n} : {data['title'][:55]}")


def make_persona(prs, data, n):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, C_BG_DARK)
    add_rect(slide, 0, 0, Inches(0.08), H, C_ACCENT)
    add_text(slide, data["title"], Inches(0.3), Inches(0.25), Inches(12.7), Inches(0.6),
             font_size=24, bold=True, color=C_WHITE)
    add_rect(slide, Inches(0.3), Inches(0.88), Inches(12.7), Inches(0.04), C_ACCENT)

    personas = data["personas"]
    col_w = Inches(4.0)
    for i, p in enumerate(personas):
        x = Inches(0.35 + i * 4.35)
        color = p["color"]

        # En-tête colonne coloré
        add_rect(slide, x, Inches(1.0), col_w, Inches(0.5), color)
        add_text(slide, f"{p['icon']}  {p['title']}", x + Inches(0.1), Inches(1.0),
                 col_w - Inches(0.1), Inches(0.5),
                 font_size=11, bold=True, color=C_BG_DARK)

        # Points
        for j, point in enumerate(p["points"]):
            y_pt = Inches(1.6 + j * 1.12)
            add_rect(slide, x, y_pt, col_w, Inches(1.0), C_BG_CARD2)
            # Bordure gauche colorée
            add_rect(slide, x, y_pt, Inches(0.05), Inches(1.0), color)
            add_text(slide, point, x + Inches(0.12), y_pt + Inches(0.12),
                     col_w - Inches(0.2), Inches(0.78),
                     font_size=10, color=C_WHITE)

    slide_n_label(slide, n)
    print(f"  ✅  Slide {n} : {data['title'][:55]}")


def make_plain(prs, data, n):
    """Slide à deux colonnes de bullets avec optionnel CTA bas."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, C_BG_DARK)
    cat_color = data.get("cat_color", C_ACCENT)
    header_bar(slide, data["title"], data.get("category"), cat_color)

    # Description
    add_text(slide, data.get("description", ""), Inches(0.3), Inches(1.3),
             Inches(12.7), Inches(0.5), font_size=10, color=C_MUTED, italic=True)

    # Colonne gauche
    add_rect(slide, Inches(0.3), Inches(1.9), Inches(6.1), Inches(0.28), cat_color)
    add_text(slide, "Pour votre offre", Inches(0.4), Inches(1.9), Inches(5.9), Inches(0.28),
             font_size=9, bold=True, color=C_BG_DARK)
    add_bullet_list(slide, data.get("bullets_left", []), Inches(0.3), Inches(2.25),
                    Inches(6.1), Inches(4.0), font_size=10.5, accent=cat_color)

    # Colonne droite
    add_rect(slide, Inches(6.7), Inches(1.9), Inches(6.3), Inches(0.28), C_ACCENT2)
    add_text(slide, "Pour votre business", Inches(6.8), Inches(1.9), Inches(6.1), Inches(0.28),
             font_size=9, bold=True, color=C_BG_DARK)
    add_bullet_list(slide, data.get("bullets_right", []), Inches(6.7), Inches(2.25),
                    Inches(6.3), Inches(4.0), font_size=10.5, accent=C_ACCENT2)

    # Séparateur vertical
    add_rect(slide, Inches(6.55), Inches(1.85), Inches(0.03), Inches(4.5), cat_color)

    # CTA bar bas
    if data.get("cta"):
        add_rect(slide, Inches(0.3), Inches(6.55), Inches(12.7), Inches(0.38), cat_color)
        add_text(slide, f"💡  {data['cta']}", Inches(0.45), Inches(6.55),
                 Inches(12.5), Inches(0.38), font_size=11, bold=True, color=C_BG_DARK)

    slide_n_label(slide, n)
    print(f"  ✅  Slide {n} : {data['title'][:55]}")


def make_features_grid(prs, data, n):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, C_BG_DARK)
    header_bar(slide, data["title"], data.get("category"), C_ACCENT)

    modules = data["modules"]  # list of (icon, name, tech)
    COLS = 4
    cell_w = Inches(3.15)
    cell_h = Inches(1.0)
    for i, (icon, name, tech) in enumerate(modules):
        col = i % COLS
        row = i // COLS
        x = Inches(0.28 + col * 3.26)
        y = Inches(1.35 + row * 1.08)
        add_rect(slide, x, y, cell_w, cell_h, C_BG_CARD2)
        # Bordure colorée alternée
        border_color = C_ACCENT if (i % 2 == 0) else C_ACCENT2
        add_rect(slide, x, y, Inches(0.05), cell_h, border_color)
        add_text(slide, icon, x + Inches(0.1), y + Inches(0.08), Inches(0.45), Inches(0.45),
                 font_size=18, color=C_WHITE)
        add_text(slide, name, x + Inches(0.55), y + Inches(0.08), cell_w - Inches(0.65), Inches(0.38),
                 font_size=10.5, bold=True, color=C_WHITE)
        add_text(slide, tech, x + Inches(0.55), y + Inches(0.5), cell_w - Inches(0.65), Inches(0.42),
                 font_size=7.5, color=C_MUTED)

    slide_n_label(slide, n)
    print(f"  ✅  Slide {n} : {data['title'][:55]}")


def make_tech(prs, data, n):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, C_BG_DARK)
    add_rect(slide, 0, 0, Inches(0.08), H, C_ACCENT)
    add_text(slide, data["title"], Inches(0.3), Inches(0.2), Inches(12.7), Inches(0.55),
             font_size=24, bold=True, color=C_WHITE)
    add_rect(slide, Inches(0.3), Inches(0.78), Inches(12.7), Inches(0.04), C_ACCENT)

    LAYER_COLORS = [C_ACCENT, C_DIST, C_GREEN, C_DEV]
    for i, (layer_name, items) in enumerate(data["layers"]):
        col_w = Inches(3.1)
        x = Inches(0.3 + i * 3.2)
        lc = LAYER_COLORS[i % len(LAYER_COLORS)]
        add_rect(slide, x, Inches(0.86), col_w - Inches(0.1), Inches(0.36), lc)
        add_text(slide, layer_name, x + Inches(0.07), Inches(0.86),
                 col_w - Inches(0.15), Inches(0.36), font_size=11, bold=True, color=C_BG_DARK)
        for j, item in enumerate(items):
            y = Inches(1.3 + j * 0.98)
            add_rect(slide, x, y, col_w - Inches(0.1), Inches(0.88), C_BG_CARD2)
            add_rect(slide, x, y, Inches(0.04), Inches(0.88), lc)
            add_text(slide, "◆", x + Inches(0.1), y + Inches(0.08),
                     Inches(0.2), Inches(0.3), font_size=7, color=lc)
            add_text(slide, item, x + Inches(0.25), y + Inches(0.06),
                     col_w - Inches(0.38), Inches(0.72), font_size=9, color=C_WHITE)
    slide_n_label(slide, n)
    print(f"  ✅  Slide {n} : {data['title'][:55]}")


def make_pricing(prs, data, n):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, C_BG_DARK)
    header_bar(slide, data["title"], data.get("category"), C_ACCENT)

    plans = data["plans"]
    col_w = Inches(4.1)
    for i, plan in enumerate(plans):
        x = Inches(0.28 + i * 4.35)
        y_start = Inches(1.35)
        bg = C_BG_CARD2
        border = plan["color"]

        if plan["highlight"]:
            # Mise en avant
            add_rect(slide, x - Inches(0.04), y_start - Inches(0.08),
                     col_w + Inches(0.08), Inches(5.6), border)
            add_rect(slide, x, y_start + Inches(0.06),
                     col_w, Inches(5.4), bg)
        else:
            add_rect(slide, x, y_start, col_w, Inches(5.3), bg)
            add_rect(slide, x, y_start, col_w, Inches(0.06), border)

        # Nom du plan
        add_text(slide, plan["name"], x + Inches(0.15), y_start + Inches(0.15),
                 col_w - Inches(0.2), Inches(0.4),
                 font_size=16, bold=True, color=border)

        # Prix
        add_text(slide, plan["price"], x + Inches(0.15), y_start + Inches(0.55),
                 col_w - Inches(0.2), Inches(0.5),
                 font_size=22, bold=True, color=C_WHITE)

        # Subtitle volume
        add_text(slide, plan["subtitle"], x + Inches(0.15), y_start + Inches(1.05),
                 col_w - Inches(0.2), Inches(0.28),
                 font_size=9, color=C_MUTED, italic=True)

        # Idéal pour
        add_rect(slide, x + Inches(0.1), y_start + Inches(1.38),
                 col_w - Inches(0.2), Inches(0.5), rgb(0x1C, 0x24, 0x36))
        add_text(slide, f"✓ Idéal pour : {plan['ideal']}", x + Inches(0.2),
                 y_start + Inches(1.4), col_w - Inches(0.3), Inches(0.45),
                 font_size=8.5, color=border)

        # Inclus
        for j, inc in enumerate(plan["includes"]):
            yy = y_start + Inches(1.98 + j * 0.41)
            r1 = add_text(slide, "✓  " + inc, x + Inches(0.18), yy,
                          col_w - Inches(0.25), Inches(0.38),
                          font_size=9.5, color=C_WHITE)

    # Note bas
    add_rect(slide, Inches(0.3), Inches(6.65), Inches(12.7), Inches(0.5), C_BG_CARD2)
    add_text(slide, f"ℹ️  {data['note']}", Inches(0.45), Inches(6.65),
             Inches(12.5), Inches(0.5), font_size=9, color=C_MUTED, italic=True)
    slide_n_label(slide, n)
    print(f"  ✅  Slide {n} : {data['title'][:55]}")


def make_timeline(prs, data, n):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, C_BG_DARK)
    STEP_COLORS = [C_DIST, C_ACCENT, C_GREEN, C_DEV]
    header_bar(slide, data["title"], data.get("category"), C_GREEN)

    steps = data["steps"]
    step_w = Inches(3.1)
    # Ligne de connexion
    add_rect(slide, Inches(0.95), Inches(2.65), Inches(11.5), Inches(0.06), rgb(0x2A, 0x35, 0x4A))

    for i, (week, title, body) in enumerate(steps):
        x = Inches(0.3 + i * 3.26)
        color = STEP_COLORS[i % len(STEP_COLORS)]

        # Cercle numéroté
        cx = x + Inches(1.3)
        add_rect(slide, cx, Inches(2.35), Inches(0.58), Inches(0.58), color)
        add_text(slide, str(i + 1), cx, Inches(2.35), Inches(0.58), Inches(0.58),
                 font_size=16, bold=True, color=C_BG_DARK, align=PP_ALIGN.CENTER)

        # Semaine
        add_text(slide, week, x, Inches(1.35), step_w, Inches(0.32),
                 font_size=9, bold=True, color=color, align=PP_ALIGN.CENTER)

        # Titre étape
        add_text(slide, title, x, Inches(1.68), step_w, Inches(0.42),
                 font_size=11, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

        # Corps
        add_rect(slide, x + Inches(0.1), Inches(3.1), step_w - Inches(0.2), Inches(2.5), C_BG_CARD2)
        add_rect(slide, x + Inches(0.1), Inches(3.1), Inches(0.05), Inches(2.5), color)
        add_text(slide, body, x + Inches(0.2), Inches(3.18),
                 step_w - Inches(0.35), Inches(2.35),
                 font_size=9.5, color=C_WHITE)

    # Note après lancement
    add_rect(slide, Inches(0.3), Inches(5.8), Inches(12.7), Inches(0.9), rgb(0x0F, 0x17, 0x28))
    add_rect(slide, Inches(0.3), Inches(5.8), Inches(0.06), Inches(0.9), C_GREEN)
    add_text(slide, f"🚀  {data['after']}", Inches(0.5), Inches(5.85),
             Inches(12.4), Inches(0.8), font_size=10, color=C_GREEN, italic=True)
    slide_n_label(slide, n)
    print(f"  ✅  Slide {n} : {data['title'][:55]}")


def make_conclusion(prs, data, n):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, C_BG_DARK)
    add_rect(slide, 0, 0, Inches(0.08), H, C_ACCENT)
    add_rect(slide, 0, H - Inches(0.08), W, Inches(0.08), C_ACCENT)
    add_rect(slide, 0, 0, W, Inches(0.08), C_ACCENT2)

    add_text(slide, data["title"], Inches(0.4), Inches(0.22), Inches(12.5), Inches(0.7),
             font_size=22, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_rect(slide, Inches(2.5), Inches(0.95), Inches(8.3), Inches(0.05), C_ACCENT)

    pts = data["points"]
    for i, pt in enumerate(pts):
        col = i % 2
        row = i // 2
        x = Inches(0.4 + col * 6.45)
        y = Inches(1.1 + row * 1.3)
        add_rect(slide, x, y, Inches(6.15), Inches(1.15), C_BG_CARD2)
        add_rect(slide, x, y, Inches(0.05), Inches(1.15), C_ACCENT)
        add_text(slide, pt, x + Inches(0.15), y + Inches(0.18),
                 Inches(5.9), Inches(0.8), font_size=11, color=C_WHITE)

    # Contact
    add_rect(slide, Inches(0.4), Inches(6.55), Inches(12.6), Inches(0.55), C_BG_CARD2)
    add_text(slide, f"🌐  {data['url']}", Inches(0.55), Inches(6.58),
             Inches(5.5), Inches(0.45), font_size=12, bold=True, color=C_ACCENT)
    add_text(slide, f"✉️  {data['contact']}", Inches(6.5), Inches(6.58),
             Inches(6.5), Inches(0.45), font_size=11, color=C_MUTED, align=PP_ALIGN.RIGHT)
    print(f"  ✅  Slide {n} : Conclusion")


# ── MAIN ───────────────────────────────────────────────────────────────────────

DISPATCH = {
    "cover":         lambda prs, d, n: make_cover(prs, d),
    "toc":           lambda prs, d, n: make_toc(prs, d),
    "stats":         make_stats,
    "persona":       make_persona,
    "plain":         make_plain,
    "features_grid": make_features_grid,
    "tech":          make_tech,
    "pricing":       make_pricing,
    "timeline":      make_timeline,
    "conclusion":    make_conclusion,
}


def generate():
    print("\n" + "═" * 62)
    print("  ☀️   HeliaPV — Présentation Marque Blanche 2026")
    print("  Cibles : Fabricants · Distributeurs · Développeurs")
    print("═" * 62 + "\n")

    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    for i, slide_data in enumerate(SLIDES, start=1):
        fn = DISPATCH.get(slide_data["type"])
        if fn:
            fn(prs, slide_data, i)
        else:
            print(f"  ⚠️  Type inconnu : {slide_data['type']}")

    prs.save(str(OUTPUT))
    print(f"\n  ✅  Fichier généré : {OUTPUT.absolute()}")
    print(f"  📊  {len(SLIDES)} slides · Format 16:9 · 1920×1080\n")


if __name__ == "__main__":
    generate()
