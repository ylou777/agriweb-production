#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Génère HeliaPV_Pipeline_PV.pptx — Charte graphique HeliaPV STRICTE
• Fond dégradé navy #0a0e27 → navy-violet #2d1b4e (OOXML bgPr)
• Accent or #FFB700 / or clair #FFD060 (dégradé shape fill)
• Police Inter
• Cartes sombres + filet couleur phase + bordure subtile
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from lxml import etree
from pptx.oxml.ns import qn

OUTPUT = "HeliaPV_Pipeline_PV.pptx"

# ── Palette HeliaPV ──────────────────────────────────────────────────────────
C_NAVY    = RGBColor(0x0A, 0x0E, 0x27)
C_NAVY2   = RGBColor(0x1A, 0x1F, 0x3A)
C_NAVY3   = RGBColor(0x2D, 0x1B, 0x4E)
C_CARD    = RGBColor(0x14, 0x1A, 0x33)
C_CARD2   = RGBColor(0x0E, 0x12, 0x26)
C_GOLD    = RGBColor(0xFF, 0xB7, 0x00)
C_GOLD_LT = RGBColor(0xFF, 0xD0, 0x60)
C_GOLD_DK = RGBColor(0xD4, 0x88, 0x00)
C_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
C_TEXT    = RGBColor(0xE2, 0xEB, 0xFF)
C_MUTED   = RGBColor(0x6A, 0x7B, 0x9C)
C_BORDER  = RGBColor(0x1D, 0x26, 0x44)
FONT      = "Inter"

# Couleurs de phase (identiques à pipeline_pv.html)
PHASE_C = {
    1: RGBColor(0xFF, 0x8C, 0x42),   # orange   – Prospection
    2: RGBColor(0x3B, 0xC4, 0xF2),   # bleu     – Études
    3: RGBColor(0xA7, 0x8B, 0xFA),   # violet   – CRM
    4: RGBColor(0x34, 0xD3, 0x99),   # vert     – Administratif
    5: RGBColor(0xF4, 0x72, 0xB6),   # rose     – Chantier
    6: RGBColor(0xFF, 0xB7, 0x00),   # or       – MES
}

W = Inches(13.33)
H = Inches(7.5)

# Stops fond dégradé : navy → navy-mid → navy-violet → navy-mid
BG_STOPS = [
    (0,      "0a0e27"),
    (40000,  "1a1f3a"),
    (75000,  "2d1b4e"),
    (100000, "1a1f3a"),
]


def shade(c, f=0.55):
    """Assombrit une RGBColor d'un facteur f."""
    return RGBColor(int(c[0] * f), int(c[1] * f), int(c[2] * f))


# ── Contenu des 6 phases ─────────────────────────────────────────────────────
PHASES = [
    {
        "num": 1,
        "title": "Prospection & Détection de Sites",
        "tagline": "Identifier les gisements PV avant la concurrence",
        "icon": "🔍",
        "features": [
            ("Cartographie GIS multi-couches",
             "RPG parcelles agricoles, PLU/PLUi, parkings, friches, toitures via GeoServer WMS/WFS. "
             "Filtres cumulables en temps réel."),
            ("Analyse LiDAR HD IGN (COPC natif)",
             "Streaming nuage de points 3D – extraction automatique des plans toiture, pente, "
             "orientation, masques d'ombrage. Résolution 10 pts/m²."),
            ("IA Helia – classification de toiture",
             "Modèle Groq LLM classifie le matériau (tuile, ardoise, métal, bac acier…) et la "
             "faisabilité PV, puis crée le prospect automatiquement dans le CRM."),
            ("Réseau HTA/BT Enedis",
             "Visualisation aérien/souterrain, mesure distance raccordement, capacité disponible "
             "par poste – données Enedis Open Data."),
            ("Google Solar Building Insights",
             "Empreinte bâtiment pré-calculée, DNI/GHI, ombres projetées voisins, détection arbres."),
            ("Recherches sauvegardées & auto-import CRM",
             "Critères de filtrage persistants, création automatique de prospects depuis résultats carte."),
        ],
        "tags": ["GeoServer WMS/WFS", "IGN LiDAR HD COPC", "Groq LLM", "Enedis HTA/BT", "Google Solar API", "RPG · PLU · Friches"],
    },
    {
        "num": 2,
        "title": "Études & Simulations Techniques",
        "tagline": "Dimensionner avec précision, documenter aux normes",
        "icon": "📊",
        "features": [
            ("Simulation PVGIS 8 760 heures",
             "3 simulateurs EU Science Hub, enveloppes statistiques P50/P90, décomposition "
             "mensuelle/saisonnière, Performance Ratio ajustable, cache local anti-quota."),
            ("Calpinage PV & plan masse",
             "Multi-zones, comptage automatique des modules, optimisation strings DC, "
             "export plan cadastral avec superposition cadrages à l'échelle."),
            ("Schéma unifilaire NF C 15-712",
             "Configuration DC strings (série/parallèle), 3 types de raccordement, stockage batterie "
             "optionnel, symboles normalisés, export PDF signable."),
            ("Autoconsommation & tarifs Enedis",
             "6 profils de consommation (RES1, RES2, PRO1, PRO2, AGR, ENT), 5 structures tarifaires "
             "(BASE, HPHC, TEMPO 3 couleurs, EJP, C4 horosaisonnier), TRI/VAN/payback."),
            ("Base données équipements CertISolis",
             "50+ modules PV (Canadian Solar, Jinko, Q Cells, REC, Sunwatt…), 30+ onduleurs "
             "(Fronius, Huawei, SMA, ABB…), paramètres électriques Voc/Vmp/Isc/Imp/Pmax."),
            ("Plans de câblage & strings",
             "Topologie strings DC, calcul section câble (mm²), cheminement, placement "
             "coffrets/compteur, export DWG/PDF."),
        ],
        "tags": ["PVGIS API", "NF C 15-712-1/2", "CertISolis", "Enedis Data Connect", "P50/P90", "IEC 61730"],
    },
    {
        "num": 3,
        "title": "Pipeline Commercial & CRM",
        "tagline": "Piloter le portefeuille, convertir les leads",
        "icon": "💼",
        "features": [
            ("CRM multi-rôles & pipeline",
             "Roles hiérarchiques admin / directeur commercial / commercial, isolation multi-tenant, "
             "stages Prospect → Qualifié → Étude → Devis → Signé → Chantier."),
            ("Génération proposition commerciale",
             "Rapport PDF auto-rempli : prévision 8 760 h, photos satellite, calpinage, "
             "schéma électrique, ROI/payback, financement – prêt à envoyer."),
            ("KPIs & tableaux de bord",
             "Taux de conversion par étape, CA prévisionnel en MW/M€, pipeline visuel, "
             "comparaison N-1, exports PDF reconductibles."),
            ("Visite technique intégrée",
             "Checklist terrain structurée, capture photos on-site (PWA-ready), "
             "relevé cotes toiture, évaluation structurelle, archivage CRM automatique."),
            ("Agenda & relances commerciales",
             "FullCalendar partagé, rappels configurables (prospect inactif, délai dépassé), "
             "historique interactions (appels, emails, visites, docs)."),
            ("Import en masse & scoring",
             "Import CSV/Excel de listes prospects, scoring automatique 0-100 selon critères "
             "pondérables (surface, distance réseau, PLU, irradiation)."),
        ],
        "tags": ["CRM Pipeline", "Multi-tenant RBAC", "Rapport PDF auto", "FullCalendar", "KPI Dashboard", "Scoring"],
    },
    {
        "num": 4,
        "title": "Dossiers Administratifs & Permis",
        "tagline": "CERFA, permis, raccordement Enedis – automatisés",
        "icon": "📋",
        "features": [
            ("Déclaration Préalable complète – CERFA 13703",
             "DP1 plan de situation cadastral, DP2 extrait cadastral avec modules, DP3 coupe "
             "transversale, DP4/DP5 façades actuelles/projetées, DP6 insertion paysagère, "
             "DP7/DP8 photos environnement."),
            ("CERFA Enedis 16702-01 pré-rempli",
             "Demande de raccordement, classification autoconsommation/injection, acheminement "
             "câbles, protection, champs signature installateur/maître d'ouvrage."),
            ("Coordonnées cadastrales automatiques",
             "Transformation IGN API cadastre → Lambert 93, extraction emprise parcelle, "
             "photomontage satellite avec superposition modules."),
            ("GED – Gestion documentaire",
             "Versioning documents, prévisualisation PDF, export ZIP dossier complet, "
             "OCR pour archivage contrats, partage lien sécurisé."),
            ("DICT, assurances & PPSPS",
             "Déclaration DT-DICT réseaux, TRC, RC chantier, plan de prévention PPSPS "
             "pour chantiers >600 m², RICT si applicable."),
            ("Suivi des délais réglementaires",
             "Jalons DP/PC avec alertes automatiques, suivi accord Enedis RACCORDEMENT, "
             "PPA si applicable, délais instruction préfecture."),
        ],
        "tags": ["CERFA 13703*09", "CERFA 16702-01", "IGN Cadastre API", "GED versionnée", "DICT/PPSPS", "Lambert 93"],
    },
    {
        "num": 5,
        "title": "Suivi de Chantier",
        "tagline": "7 sous-phases normalisées, de l'études à la livraison",
        "icon": "🏗️",
        "features": [
            ("ENG – Études & Ingénierie (30 j)",
             "Validation PVsyst P50/P90, note de calcul DC/AC, plans IFC d'exécution, "
             "CCTP/DPGF/DQE, analyse d'ombrage, sélection matériel définitive."),
            ("ADM – Administratif & Permitting (90 j)",
             "DICT réseaux, dépôt DP/PC, accord raccordement Enedis, PPA si applicable, "
             "assurances TRC + RC chantier, PPSPS, RICT."),
            ("APPRO – Approvisionnement (45 j)",
             "BdC modules/onduleurs/structure, réception avec flashlisting (courbes I-V), "
             "tracking numéros de série en base, rapport de réception signé."),
            ("GC – Génie Civil & Charpente (20 j)",
             "Préparation toiture, renfort chevrons, étanchéité, pose rails, nivellement, "
             "anti-corrosion, procès-verbal de réception levage."),
            ("INST – Installation Électrique (30 j)",
             "Pose modules sur rails, tests strings DC, coffrets AC, disjoncteur principal, "
             "mise à la terre, cheminement et terminaison câbles."),
            ("COMMI – Mise en Service & Tests (10 j)",
             "Mégohmmétrie 1 000 V DC, test polarité, impédance boucle AC, "
             "anti-îlotage, analyse harmoniques, rapport IEC 62446-1 complet."),
        ],
        "tags": ["IEC 62446-1", "PPSPS/RICT", "NCF/DOE", "Flashlisting I-V", "Mégohmètre 1kV", "PVsyst P50/P90"],
    },
    {
        "num": 6,
        "title": "Livraison & Mise en Service",
        "tagline": "DOE, CONSUEL, monitoring activé, dossier bancaire",
        "icon": "⚡",
        "features": [
            ("DOE & réception préfecture",
             "Déclaration de fin de travaux, dépôt préfecture, rapport IEC 62446-1 complet "
             "généré automatiquement avec tous les PV de tests."),
            ("Inspection finale Enedis & CONSUEL",
             "Dossier technique de consignation S21, attestation de conformité CONSUEL, "
             "mise en service réseau, contrat d'accès."),
            ("Rapport IEC 62446 & activation monitoring",
             "Rapport de mise en service normé, activation système de monitoring, "
             "comparaison production réelle vs PVGIS (P50), alertes de sous-performance."),
            ("Documentation bancaire – Lender Pack",
             "Dossier Due Diligence complet pour financement/refinancement, contrat PPA, "
             "rapports de performance investisseurs, valorisation actif."),
            ("Remise O&M & garanties fabricants",
             "Manuels d'exploitation, enregistrement garanties (25 ans modules, 10 ans "
             "onduleurs), mise en place contrat O&M préventif."),
            ("Blog technique & veille réglementaire",
             "Articles techniques (LiDAR IGN, PVGIS, Enedis Data Connect, PLU/GPU), "
             "actualisation automatique barèmes CRE, landing pages bureaux d'études."),
        ],
        "tags": ["IEC 62446", "CONSUEL", "DOE Préfecture", "Monitoring", "Lender Pack Due Diligence", "O&M 25 ans"],
    },
]

# ── Helpers ──────────────────────────────────────────────────────────────────

def set_gradient_bg(slide, stops=None, angle=8100000):
    """Fond dégradé via manipulation XML OOXML (bgPr). angle=8100000 → 135°."""
    if stops is None:
        stops = BG_STOPS
    sld = slide._element
    cSld = sld.find(qn('p:cSld'))
    if cSld is None:
        return
    for b in cSld.findall(qn('p:bg')):
        cSld.remove(b)
    bg = etree.Element(qn('p:bg'))
    bgPr = etree.SubElement(bg, qn('p:bgPr'))
    gf = etree.SubElement(bgPr, qn('a:gradFill'))
    gsLst = etree.SubElement(gf, qn('a:gsLst'))
    for pos, col in stops:
        gs = etree.SubElement(gsLst, qn('a:gs'))
        gs.set('pos', str(int(pos)))
        sc = etree.SubElement(gs, qn('a:srgbClr'))
        sc.set('val', col)
    lin = etree.SubElement(gf, qn('a:lin'))
    lin.set('ang', str(angle))
    lin.set('scaled', '0')
    etree.SubElement(bgPr, qn('a:effectLst'))
    cSld.insert(0, bg)


def solid_rect(slide, l, t, w, h, fill, border=None, border_w=Pt(0.75)):
    s = slide.shapes.add_shape(1, int(l), int(t), int(w), int(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if border:
        s.line.color.rgb = border
        s.line.width = int(border_w)
    else:
        s.line.fill.background()
    return s


def grad_rect(slide, l, t, w, h, c1, c2, angle=90.0, border=None):
    """Rectangle dégradé 2 stops (API python-pptx)."""
    s = slide.shapes.add_shape(1, int(l), int(t), int(w), int(h))
    f = s.fill
    f.gradient()
    f.gradient_angle = angle
    f.gradient_stops[0].position = 0.0
    f.gradient_stops[0].color.rgb = c1
    f.gradient_stops[1].position = 1.0
    f.gradient_stops[1].color.rgb = c2
    if border:
        s.line.color.rgb = border
    else:
        s.line.fill.background()
    return s


def add_txt(slide, text, l, t, w, h, size=12, bold=False, italic=False,
            color=None, align=PP_ALIGN.LEFT):
    if color is None:
        color = C_TEXT
    tb = slide.shapes.add_textbox(int(l), int(t), int(w), int(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = FONT
    return tb


def gold_line(slide, top):
    grad_rect(slide, Inches(0.4), top, Inches(12.5), Pt(2),
              C_GOLD_LT, C_GOLD_DK, angle=0.0)


def chrome(slide, n=None):
    """Filet vertical or gauche + filet or bas + numéro slide."""
    grad_rect(slide, 0, 0, Pt(5), H, C_GOLD, C_GOLD_DK, angle=90.0)
    grad_rect(slide, 0, H - Pt(5), W, Pt(5), C_GOLD_DK, C_GOLD, angle=0.0)
    if n:
        add_txt(slide, str(n), W - Inches(0.45), H - Inches(0.38),
                Inches(0.35), Inches(0.28), size=9,
                color=RGBColor(0x30, 0x3C, 0x5A), align=PP_ALIGN.RIGHT)


# ── Slide 1 — Couverture ─────────────────────────────────────────────────────

def slide_cover(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_gradient_bg(sl)
    chrome(sl, 1)

    # Panneau droit violet
    solid_rect(sl, W - Inches(5.4), 0, Inches(5.4), H, C_NAVY3)

    # Cercle solaire (anneau or)
    sun = sl.shapes.add_shape(1,
        int(W - Inches(5.0)), int(Inches(0.35)),
        int(Inches(4.2)), int(Inches(4.2)))
    sun.fill.solid(); sun.fill.fore_color.rgb = C_GOLD
    sun.line.fill.background()
    sun_h = sl.shapes.add_shape(1,
        int(W - Inches(4.78)), int(Inches(0.57)),
        int(Inches(3.76)), int(Inches(3.76)))
    sun_h.fill.solid(); sun_h.fill.fore_color.rgb = C_NAVY3
    sun_h.line.fill.background()
    sun_core = sl.shapes.add_shape(1,
        int(W - Inches(3.6)), int(Inches(1.75)),
        int(Inches(1.42)), int(Inches(1.42)))
    sun_core.fill.solid(); sun_core.fill.fore_color.rgb = C_GOLD_LT
    sun_core.line.fill.background()

    # Badge HeliaPV
    solid_rect(sl, Inches(0.5), Inches(0.82), Inches(3.5), Inches(0.34),
               RGBColor(0x18, 0x1E, 0x3C), border=RGBColor(0x2A, 0x36, 0x5A))
    add_txt(sl, "\u2600  HeliaPV  \u00b7  Solution Int\u00e9gr\u00e9e",
            Inches(0.58), Inches(0.83), Inches(3.35), Inches(0.32),
            size=9, bold=True, color=C_GOLD)

    # Titre
    add_txt(sl, "Du Terrain \u00e0 la Mise en Service",
            Inches(0.5), Inches(1.42), Inches(8.2), Inches(0.72),
            size=34, bold=True, color=C_WHITE)
    add_txt(sl, "\u2014 Une Seule Plateforme",
            Inches(0.5), Inches(2.12), Inches(8.2), Inches(0.52),
            size=22, bold=True, color=C_GOLD)

    # Sous-titre
    add_txt(sl, (
        "HeliaPV couvre la totalit\u00e9 du cycle de vie d\u2019un projet "
        "photovolta\u00efque\u00a0: de la prospection cartographique jusqu\u2019au "
        "suivi de chantier et la mise en service, en passant par les \u00e9tudes "
        "techniques, le pipeline CRM et les dossiers administratifs."
    ), Inches(0.5), Inches(2.82), Inches(7.8), Inches(1.0),
        size=10.5, color=C_MUTED, italic=True)

    # Pipeline 6 phases (bandes colorées)
    phases_lbl = ["01 Prospection", "02 \u00c9tudes",
                  "03 CRM", "04 Administratif", "05 Chantier", "06 MES"]
    for i, (lbl, pc) in enumerate(zip(phases_lbl, PHASE_C.values())):
        x = Inches(0.5 + i * 2.05)
        grad_rect(sl, x, Inches(4.12), Inches(1.95), Inches(0.38),
                  pc, shade(pc, 0.6), angle=90.0)
        add_txt(sl, lbl, x, Inches(4.13), Inches(1.95), Inches(0.36),
                size=8.5, bold=True, color=C_NAVY, align=PP_ALIGN.CENTER)

    # Ligne séparatrice
    grad_rect(sl, Inches(0.5), Inches(4.62), Inches(12.3), Pt(1),
              C_BORDER, C_BORDER)
    add_txt(sl,
            "8\u202f760\u00a0h PVGIS  \u00b7  LiDAR IGN COPC  \u00b7  NF\u00a0C\u00a015-712  "
            "\u00b7  IEC\u00a062446  \u00b7  CERFA\u00a016702  \u00b7  Groq LLM",
            Inches(0.5), Inches(4.68), Inches(12.3), Inches(0.3),
            size=8.5, color=RGBColor(0x44, 0x52, 0x72),
            italic=True, align=PP_ALIGN.CENTER)

    # Version / URL
    add_txt(sl, "Version 3.1  \u00b7  Mars 2026",
            Inches(0.5), Inches(6.9), Inches(4), Inches(0.3),
            size=8.5, color=RGBColor(0x3A, 0x46, 0x66))
    add_txt(sl, "app.heliapv.fr",
            W - Inches(2.6), Inches(6.9), Inches(2.1), Inches(0.3),
            size=9, bold=True, color=C_GOLD, align=PP_ALIGN.RIGHT)

    print("  \u2705  Slide 1 : Couverture")


# ── Slide 2 — Vue d'ensemble ─────────────────────────────────────────────────

def slide_overview(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_gradient_bg(sl)
    chrome(sl, 2)

    add_txt(sl, "Vue d\u2019ensemble \u2014 Pipeline Projet PV Complet",
            Inches(0.4), Inches(0.18), Inches(12.5), Inches(0.6),
            size=22, bold=True, color=C_WHITE)
    gold_line(sl, Inches(0.82))

    labels = ["Prospection\n& D\u00e9tection", "\u00c9tudes &\nSimulations",
              "Pipeline\nCommercial", "Dossiers\nAdministratifs",
              "Suivi\nChantier", "Livraison\n& MES"]
    icons = ["\U0001F50D", "\U0001F4CA", "\U0001F4BC",
             "\U0001F4CB", "\U0001F3D7", "\u26A1"]
    details = [
        "Cartographie GIS\nLiDAR HD IGN COPC\nIA Helia (Groq)\nEnedis HTA/BT",
        "PVGIS 8\u202f760\u00a0h\nCalpinage + Masse\nSch\u00e9ma unifilaire\nAutoconsommation",
        "CRM multi-r\u00f4les\nProposition PDF\nKPIs & Agenda\nVisite technique",
        "DP CERFA 13703\nCERFA Enedis 16702\nGED documentaire\nDICT / PPSPS",
        "ENG \u00b7 ADM \u00b7 APPRO\nGC \u00b7 INST \u00b7 COMMI\nFlashlisting I-V\nTests IEC 62446",
        "DOE Pr\u00e9fecture\nCONSUEL\nMonitoring\nLender Pack",
    ]

    col_w = Inches(2.09)
    for i, (label, icon, detail, pc) in enumerate(
            zip(labels, icons, details, PHASE_C.values())):
        x = Inches(0.25 + i * 2.17)
        # En-tête dégradé couleur phase
        grad_rect(sl, x, Inches(1.05), col_w, Inches(0.44),
                  pc, shade(pc, 0.6), angle=90.0)
        add_txt(sl, f"{icon}  0{i+1}", x, Inches(1.06), col_w, Inches(0.42),
                size=11.5, bold=True, color=C_NAVY, align=PP_ALIGN.CENTER)
        # Titre phase
        solid_rect(sl, x, Inches(1.49), col_w, Inches(0.58), C_CARD)
        add_txt(sl, label, x + Inches(0.06), Inches(1.50),
                col_w - Inches(0.08), Inches(0.56),
                size=9, bold=True, color=pc)
        # Corps
        solid_rect(sl, x, Inches(2.07), col_w, Inches(4.64),
                   C_CARD2, border=C_BORDER)
        add_txt(sl, detail,
                x + Inches(0.12), Inches(2.18),
                col_w - Inches(0.16), Inches(4.4),
                size=9, color=RGBColor(0xB0, 0xC0, 0xDC))
        # Flèche
        if i < 5:
            add_txt(sl, "\u203a",
                    x + col_w - Inches(0.04), Inches(4.1),
                    Inches(0.22), Inches(0.38),
                    size=18, bold=True,
                    color=RGBColor(0x28, 0x34, 0x52),
                    align=PP_ALIGN.CENTER)

    add_txt(sl,
            "Chaque phase est enti\u00e8rement int\u00e9gr\u00e9e dans HeliaPV \u2014 aucun outil tiers requis.",
            Inches(0.4), Inches(7.05), Inches(12.5), Inches(0.28),
            size=8.5, color=C_MUTED, italic=True, align=PP_ALIGN.CENTER)

    print("  \u2705  Slide 2 : Vue d\u2019ensemble")


# ── Slide 3 — Stats clés ─────────────────────────────────────────────────────

def slide_stats(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_gradient_bg(sl)
    chrome(sl, 3)

    add_txt(sl, "Chiffres Cl\u00e9s \u2014 Ce que couvre HeliaPV",
            Inches(0.4), Inches(0.18), Inches(12.5), Inches(0.6),
            size=22, bold=True, color=C_WHITE)
    gold_line(sl, Inches(0.82))

    data = [
        ("6",         "Phases projet\ncouvertes A\u2192Z",          PHASE_C[1]),
        ("8\u202f760\u00a0h", "Simulation PVGIS\nheure par heure",  PHASE_C[2]),
        ("80+",       "\u00c9quipements\nBDD CertISolis",            PHASE_C[3]),
        ("7",         "Sous-phases chantier\nnormalis\u00e9es IEC",  PHASE_C[4]),
        ("10\u00a0pts/m\u00b2", "LiDAR HD IGN COPC\nr\u00e9solution toiture", PHASE_C[5]),
        ("2 CERFA",   "16702 Enedis\u00a0+\n13703 DP auto",          PHASE_C[6]),
        ("5 tarifs",  "Enedis BASE/HPHC\nTEMPO/EJP/C4",             PHASE_C[1]),
        ("NF\u00a0+\u00a0IEC", "C15-712 \u00b7 62446-1\nsch\u00e9mas norm\u00e9s", PHASE_C[2]),
        ("Groq LLM",  "IA Helia int\u00e9gr\u00e9e\nclassif. toiture", PHASE_C[3]),
    ]

    for i, (val, label, color) in enumerate(data):
        col = i % 3
        row = i // 3
        x = Inches(0.4 + col * 4.28)
        y = Inches(1.05 + row * 1.95)
        solid_rect(sl, x, y, Inches(4.0), Inches(1.85), C_CARD, border=C_BORDER)
        grad_rect(sl, x, y, Pt(4), Inches(1.85),
                  color, shade(color, 0.6), angle=90.0)
        add_txt(sl, val, x + Inches(0.18), y + Inches(0.16),
                Inches(3.75), Inches(0.82), size=30, bold=True, color=color)
        add_txt(sl, label, x + Inches(0.18), y + Inches(1.02),
                Inches(3.75), Inches(0.75),
                size=10, color=RGBColor(0xB0, 0xC0, 0xDC))

    print("  \u2705  Slide 3 : Stats cl\u00e9s")


# ── Slides 4–9 — Une slide par phase ─────────────────────────────────────────

def slide_phase(prs, phase, slide_n):
    n  = phase["num"]
    pc = PHASE_C[n]
    dk = shade(pc, 0.55)

    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_gradient_bg(sl)
    chrome(sl, slide_n)

    # Bandeau titre — dégradé navy2 → navy3 (cohérent pour toutes les phases)
    grad_rect(sl, Pt(5), 0, W - Pt(10), Inches(1.08),
              C_NAVY2, C_NAVY3, angle=0.0)
    # Filet couleur phase sous le bandeau
    grad_rect(sl, Pt(5), Inches(1.05), W - Pt(10), Pt(3),
              pc, dk, angle=0.0)

    # Carré numéro phase (dégradé couleur phase)
    grad_rect(sl, Inches(0.22), Inches(0.1), Inches(0.86), Inches(0.86),
              pc, dk, angle=90.0)
    add_txt(sl, f"0{n}", Inches(0.22), Inches(0.1), Inches(0.86), Inches(0.5),
            size=20, bold=True, color=C_NAVY, align=PP_ALIGN.CENTER)
    add_txt(sl, phase["icon"],
            Inches(0.22), Inches(0.56), Inches(0.86), Inches(0.38),
            size=14, color=C_NAVY, align=PP_ALIGN.CENTER)

    # Titre + tagline
    add_txt(sl, phase["title"],
            Inches(1.22), Inches(0.1), Inches(9.8), Inches(0.54),
            size=21, bold=True, color=C_WHITE)
    add_txt(sl, phase["tagline"],
            Inches(1.22), Inches(0.64), Inches(9.8), Inches(0.32),
            size=10.5, italic=True, color=pc)

    # Tags (haut droit)
    tx = W - Inches(3.5)
    ty = Inches(0.12)
    for tag in phase["tags"][:4]:
        tw = min(Inches(3.2), Pt(len(tag) * 5.5) + Pt(20))
        solid_rect(sl, tx, ty, tw, Inches(0.24),
                   RGBColor(0x14, 0x1A, 0x32),
                   border=RGBColor(int(pc[0] * 0.3), int(pc[1] * 0.3), int(pc[2] * 0.3)))
        add_txt(sl, f"\u2022 {tag}",
                tx + Inches(0.08), ty + Inches(0.01),
                tw - Inches(0.1), Inches(0.22),
                size=7.5, bold=True, color=pc)
        ty += Inches(0.27)

    # 6 feature cards — 2 colonnes × 3 lignes
    for i, (feat_title, feat_desc) in enumerate(phase["features"]):
        col = i % 2
        row = i // 2
        fx = Inches(0.16 + col * 6.6)
        fy = Inches(1.15 + row * 1.78)
        fw = Inches(6.35)
        fh = Inches(1.72)

        solid_rect(sl, fx, fy, fw, fh, C_CARD, border=C_BORDER)
        grad_rect(sl, fx, fy, Pt(4), fh, pc, dk, angle=90.0)

        add_txt(sl, feat_title,
                fx + Inches(0.14), fy + Inches(0.1),
                fw - Inches(0.2), Inches(0.32),
                size=10.5, bold=True, color=C_WHITE)
        add_txt(sl, feat_desc,
                fx + Inches(0.14), fy + Inches(0.43),
                fw - Inches(0.22), Inches(1.18),
                size=8.8, color=RGBColor(0xA0, 0xB2, 0xCC))

    # Tags bas de slide
    tx = Inches(0.16)
    ty = Inches(6.72)
    for tag in phase["tags"]:
        tw = min(Inches(3.0), Pt(len(tag) * 5.5) + Pt(22))
        if tx + tw > W - Inches(0.3):
            break
        solid_rect(sl, tx, ty, tw, Inches(0.26),
                   RGBColor(0x16, 0x1C, 0x32),
                   border=RGBColor(int(pc[0] * 0.22), int(pc[1] * 0.22), int(pc[2] * 0.22)))
        add_txt(sl, tag, tx + Inches(0.08), ty + Inches(0.02),
                tw - Inches(0.1), Inches(0.22),
                size=7.5, bold=True, color=pc)
        tx += tw + Inches(0.1)

    print(f"  \u2705  Slide {slide_n} : Phase {n} \u2014 {phase['title'][:45]}")


# ── Slide 10 — Architecture technique ────────────────────────────────────────

def slide_architecture(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_gradient_bg(sl)
    chrome(sl, 10)

    add_txt(sl, "Architecture Technique",
            Inches(0.4), Inches(0.18), Inches(12.5), Inches(0.6),
            size=22, bold=True, color=C_WHITE)
    gold_line(sl, Inches(0.82))

    layers = [
        ("\U0001F5FA  Frontend",       PHASE_C[1], [
            "Leaflet.js 1.9 \u2014 carte interactive multi-couches",
            "Three.js / WebGL \u2014 LiDAR 3D interactif",
            "Vanilla JS + Bootstrap 5 \u2014 UI composants",
            "Chart.js \u2014 graphiques analytiques",
            "FullCalendar.js \u2014 agenda CRM",
            "SSE EventSource \u2014 streaming temps r\u00e9el",
        ]),
        ("\u2699\ufe0f  Backend",       PHASE_C[2], [
            "Python 3.11 \u00b7 Flask 3.1",
            "GeoPandas 0.14 \u00b7 Shapely 2.0",
            "laspy[lazrs] \u2014 COPC LiDAR streaming",
            "reportlab + python-docx \u2014 g\u00e9n. docs",
            "Groq SDK \u2014 Helia IA",
            "python-pptx \u00b7 jsPDF \u2014 exports",
        ]),
        ("\U0001F6F0  APIs & Data",     PHASE_C[4], [
            "IGN G\u00e9oportail \u2014 cadastre, ortho, LiDAR",
            "PVGIS JRC \u2014 irradiation 8\u202f760\u00a0h",
            "API BAN \u2014 g\u00e9ocodage fran\u00e7ais",
            "G\u00e9oRisques \u2014 risques naturels",
            "Enedis Open Data \u2014 HTA/BT OAuth2",
            "Google Solar API \u2014 Building Insights",
        ]),
        ("\u2601\ufe0f  Infrastructure", PHASE_C[6], [
            "Railway.app \u2014 cloud PaaS Docker",
            "PostgreSQL \u2014 base production",
            "GeoServer \u2014 WMS/WFS couches g\u00e9o",
            "Gunicorn \u2014 WSGI production",
            "GitHub CI/CD \u2014 d\u00e9ploiement auto",
            "PBKDF2 \u00b7 HTTPS forc\u00e9 \u00b7 Flask-Limiter",
        ]),
    ]

    col_w = Inches(3.14)
    for i, (name, color, items) in enumerate(layers):
        x = Inches(0.22 + i * 3.26)
        dk = shade(color, 0.6)
        grad_rect(sl, x, Inches(1.0), col_w, Inches(0.46),
                  color, dk, angle=90.0)
        add_txt(sl, name, x + Inches(0.1), Inches(1.02),
                col_w - Inches(0.12), Inches(0.42),
                size=10, bold=True, color=C_NAVY)
        for j, item in enumerate(items):
            y = Inches(1.52 + j * 0.84)
            bg = C_CARD if j % 2 == 0 else C_CARD2
            solid_rect(sl, x, y, col_w, Inches(0.80), bg, border=C_BORDER)
            grad_rect(sl, x, y, Pt(4), Inches(0.80), color, dk, angle=90.0)
            add_txt(sl, item, x + Inches(0.12), y + Inches(0.12),
                    col_w - Inches(0.18), Inches(0.58),
                    size=8.5, color=RGBColor(0xB0, 0xC0, 0xDC))

    print("  \u2705  Slide 10 : Architecture technique")


# ── Slide 11 — Conclusion ────────────────────────────────────────────────────

def slide_conclusion(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_gradient_bg(sl)
    chrome(sl, 11)

    # Fond droit violet foncé
    solid_rect(sl, Inches(7.8), 0, W - Inches(7.8), H, C_NAVY3)

    # Cercle solaire
    sun = sl.shapes.add_shape(1,
        int(W - Inches(5.1)), int(Inches(0.6)),
        int(Inches(4.4)), int(Inches(4.4)))
    sun.fill.solid(); sun.fill.fore_color.rgb = C_GOLD
    sun.line.fill.background()
    sun_h = sl.shapes.add_shape(1,
        int(W - Inches(4.86)), int(Inches(0.84)),
        int(Inches(3.92)), int(Inches(3.92)))
    sun_h.fill.solid(); sun_h.fill.fore_color.rgb = C_NAVY3
    sun_h.line.fill.background()
    sun_core = sl.shapes.add_shape(1,
        int(W - Inches(3.6)), int(Inches(2.1)),
        int(Inches(1.45)), int(Inches(1.45)))
    sun_core.fill.solid(); sun_core.fill.fore_color.rgb = C_GOLD_LT
    sun_core.line.fill.background()

    # Badge
    solid_rect(sl, Inches(0.5), Inches(0.72), Inches(3.4), Inches(0.34),
               RGBColor(0x18, 0x1E, 0x3C), border=RGBColor(0x2A, 0x36, 0x5A))
    add_txt(sl, "\u2600  HeliaPV  \u00b7  Solution Compl\u00e8te",
            Inches(0.58), Inches(0.73), Inches(3.28), Inches(0.32),
            size=8.5, bold=True, color=C_GOLD)

    add_txt(sl, "La plateforme compl\u00e8te",
            Inches(0.5), Inches(1.26), Inches(7.2), Inches(0.62),
            size=26, bold=True, color=C_WHITE)
    add_txt(sl, "du d\u00e9veloppeur solaire",
            Inches(0.5), Inches(1.84), Inches(7.2), Inches(0.62),
            size=26, bold=True, color=C_GOLD)

    points = [
        (PHASE_C[1], "Cartographie GIS + LiDAR HD IGN COPC \u2014 extraction plans toiture 3D"),
        (PHASE_C[2], "Simulation PVGIS 8\u202f760\u00a0h P50/P90 + calpinage + sch\u00e9ma NF\u00a0C\u00a015-712"),
        (PHASE_C[3], "CRM int\u00e9gr\u00e9 multi-r\u00f4les du premier contact \u00e0 la mise en service"),
        (PHASE_C[4], "CERFA 13703 + 16702 g\u00e9n\u00e9r\u00e9s auto + GED versionn\u00e9e + PPSPS"),
        (PHASE_C[5], "Suivi chantier 7 sous-phases IEC\u00a062446-1 \u2014 m\u00e9gohmm\u00e8tre, anti-\u00eelotage"),
        (PHASE_C[6], "DOE pr\u00e9fecture + CONSUEL + monitoring + lender pack bancaire"),
        (C_GOLD,     "IA Helia (Groq LLM) sp\u00e9cialis\u00e9e solaire, int\u00e9gr\u00e9e nativement"),
        (C_MUTED,    "SaaS cloud-native Railway \u00b7 Docker \u00b7 CI/CD GitHub"),
    ]

    for i, (color, point) in enumerate(points):
        y = Inches(2.72 + i * 0.545)
        grad_rect(sl, Inches(0.5), y, Pt(4.5), Inches(0.38),
                  color, shade(color, 0.6), angle=90.0)
        add_txt(sl, point, Inches(1.0), y + Inches(0.06),
                Inches(6.7), Inches(0.36), size=9.5,
                color=RGBColor(0xC2, 0xCE, 0xE8))

    # CTA
    grad_rect(sl, Inches(0.5), Inches(7.04), Inches(3.2), Inches(0.38),
              C_GOLD_LT, C_GOLD_DK, angle=0.0)
    add_txt(sl, "app.heliapv.fr",
            Inches(0.5), Inches(7.04), Inches(3.2), Inches(0.38),
            size=13, bold=True, color=C_NAVY, align=PP_ALIGN.CENTER)
    add_txt(sl, "contact@heliapv.fr",
            Inches(3.88), Inches(7.09), Inches(3.5), Inches(0.3),
            size=9.5, color=C_MUTED)

    print("  \u2705  Slide 11 : Conclusion")


# ── Main ─────────────────────────────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    print("\n\U0001F680  G\u00e9n\u00e9ration HeliaPV_Pipeline_PV.pptx (charte HeliaPV stricte)\n")

    slide_cover(prs)
    slide_overview(prs)
    slide_stats(prs)
    for i, phase in enumerate(PHASES):
        slide_phase(prs, phase, 4 + i)
    slide_architecture(prs)
    slide_conclusion(prs)

    prs.save(OUTPUT)
    print(f"\n\u2705  Fichier g\u00e9n\u00e9r\u00e9 : {OUTPUT}")
    print(f"   {len(prs.slides)} slides \u00b7 Charte HeliaPV respect\u00e9e\n")


if __name__ == "__main__":
    main()
