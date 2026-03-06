#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Génère une présentation PowerPoint complète d'AgriWeb / HeliaPV
à partir des screenshots capturés par site_explorer_agent.py
"""

import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.oxml.ns as nsmap
from lxml import etree

# ── Répertoires ──────────────────────────────────────────────────────────────
SCREENSHOTS = Path("site_explorer_output/screenshots")
OUTPUT      = Path("site_explorer_output/HeliaPV_Presentation.pptx")

# ── Palette couleurs HeliaPV ──────────────────────────────────────────────────
C_BG_DARK   = RGBColor(0x0A, 0x0E, 0x1A)   # fond très sombre
C_BG_CARD   = RGBColor(0x11, 0x18, 0x27)   # surface cartes
C_ACCENT    = RGBColor(0xF5, 0x9E, 0x0B)   # jaune solaire
C_ACCENT2   = RGBColor(0x63, 0x66, 0xF1)   # indigo
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
C_MUTED     = RGBColor(0x94, 0xA3, 0xB8)
C_OK        = RGBColor(0x10, 0xB9, 0x81)   # vert

# ── Dimensions slide 16:9 ─────────────────────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)

# ─────────────────────────────────────────────────────────────────────────────
# CONTENU TECHNIQUE DE CHAQUE SLIDE
# ─────────────────────────────────────────────────────────────────────────────

SLIDES = [

    # ── SLIDE 1 : Couverture ─────────────────────────────────────────────────
    {
        "type": "cover",
        "title": "HeliaPV / AgriWeb",
        "subtitle": "Plateforme SaaS de prospection et développement de projets photovoltaïques",
        "version": "Version 3.2  ·  Mars 2026",
        "tagline": "Cartographie intelligente · CRM intégré · IA Solaire · Analyses LiDAR HD",
    },

    # ── SLIDE 2 : Sommaire ───────────────────────────────────────────────────
    {
        "type": "toc",
        "title": "Sommaire",
        "sections": [
            ("1", "Présentation générale & architecture",           "Page d'accueil, stack technique"),
            ("2", "Authentification & gestion des accès",           "Login sécurisé, rôles, sessions"),
            ("3", "Application cartographique principale",          "Carte interactive, couches de données"),
            ("4", "Recherche & découverte de sites",                "Adresse, commune, filtres avancés"),
            ("5", "CRM – Gestion commerciale",                      "Pipeline, dashboard, statistiques"),
            ("6", "Calpinage photovoltaïque 3D",                    "Simulation PVGIS, rendement, autoconso"),
            ("7", "Rapports & analyses techniques",                 "Rapport commune, point GPS, PLU, GeoRisques"),
            ("8", "Modules spécialisés",                            "LiDAR HD, AO PV, Visite technique, Chantier"),
            ("9", "Intelligence Artificielle – Helia",              "Assistant solaire IA, analyse projets"),
            ("10", "Administration & déploiement",                  "Backoffice, Railway, scalabilité"),
        ],
    },

    # ── SLIDE 3 : Page d'accueil ─────────────────────────────────────────────
    {
        "type": "feature",
        "category": "PRÉSENTATION GÉNÉRALE",
        "title": "Page d'accueil – Landing Page",
        "screenshot": "Page_d_accueil.png",
        "description": (
            "Interface d'entrée de la plateforme HeliaPV avec présentation des fonctionnalités clés "
            "et accès direct à l'application."
        ),
        "bullets": [
            "Design responsive dark-theme optimisé conversion",
            "Accès direct à la démo sans inscription",
            "Présentation des modules : Cartographie · CRM · IA · LiDAR",
            "Charte graphique solaire : gradients doré/indigo sur fond sombre",
            "Intégration Stripe pour souscription en ligne",
            "Compatible mobile / tablette / desktop 4K",
        ],
        "tech": "Flask · Jinja2 · TailwindCSS · Leaflet.js",
    },

    # ── SLIDE 4 : Authentification ───────────────────────────────────────────
    {
        "type": "feature",
        "category": "AUTHENTIFICATION & SÉCURITÉ",
        "title": "Système d'authentification sécurisé",
        "screenshot": "Formulaire_connexion.png",
        "screenshot2": "00_login_form.png",
        "description": (
            "Système d'authentification multi-niveaux avec gestion des rôles, "
            "sessions JWT et confirmation par email."
        ),
        "bullets": [
            "Authentification par email + mot de passe hashé (bcrypt)",
            "Tokens JWT signés (PyJWT) avec expiration configurable",
            "Confirmation email obligatoire à l'inscription",
            "Gestion des rôles : admin / user / trial / premium",
            "Sessions sécurisées côté serveur (Flask-Session)",
            "Protection CSRF sur tous les formulaires POST",
            "Limitation de taux (Flask-Limiter) anti-brute-force",
            "Réinitialisation mot de passe par lien temporaire",
        ],
        "tech": "Flask-Login · bcrypt · PyJWT · Flask-Limiter · SQLite/PostgreSQL",
    },

    # ── SLIDE 5 : Carte principale ───────────────────────────────────────────
    {
        "type": "feature",
        "category": "APPLICATION CARTOGRAPHIQUE",
        "title": "Carte interactive principale",
        "screenshot": "Carte_principale.png",
        "description": (
            "Cœur de la plateforme : carte interactive multi-couches avec données géospatiales "
            "en temps réel pour identifier les opportunités photovoltaïques."
        ),
        "bullets": [
            "Carte Leaflet.js avec fonds IGN, OSM, Satellite (WMTS)",
            "Couches : RPG parcelles agricoles · Postes HTA/BT · Friches · PLU · ZAER",
            "Recherche unifiée : adresse (BAN) + commune + coordonnées GPS",
            "Panneau latéral de filtres dynamiques (surface, distance réseau, type)",
            "Export KML / GeoJSON des résultats",
            "Intégration GeoServer pour couches WMS/WFS personnalisées",
            "Affichage données cadastrales (API Géoportail IGN)",
            "Calcul potentiel solaire en survol de parcelle (irradiation PVGIS)",
        ],
        "tech": "Leaflet.js · GeoServer · WMS/WFS · IGN APIs · PVGIS · Shapely · GeoPandas",
    },

    # ── SLIDE 6 : Recherche adresse ──────────────────────────────────────────
    {
        "type": "feature",
        "category": "RECHERCHE & DÉCOUVERTE",
        "title": "Recherche par adresse – Analyse solaire instantanée",
        "screenshot": "recherche_adresse_resultats.png",
        "screenshot2": "recherche_adresse_saisie.png",
        "description": (
            "La saisie d'une adresse déclenche une analyse complète de la parcelle : "
            "données cadastrales, potentiel solaire, distances réseau, contraintes PLU."
        ),
        "bullets": [
            "Autocomplétion BAN (Base Adresse Nationale) en temps réel",
            "Géocodage instantané → centrage carte + identification parcelle cadastrale",
            "Calcul irradiation solaire via PVGIS API (kWh/kWc/an)",
            "Détection orientation et inclinaison de toiture (analyse DSM/LiDAR)",
            "Distance aux postes de transformation HTA et BT (OpenStreetMap + Enedis)",
            "Vérification contraintes urbanistiques (PLU, PPRI, zones protégées)",
            "Données propriétaires via API SIRENE (SIREN/SIRET)",
            "Affichage des risques naturels (GéoRisques API)",
        ],
        "tech": "API BAN · PVGIS · API Géoportail · OpenStreetMap · GéoRisques · API SIRENE",
    },

    # ── SLIDE 7 : Recherche commune ──────────────────────────────────────────
    {
        "type": "feature",
        "category": "RECHERCHE & DÉCOUVERTE",
        "title": "Recherche par commune – Prospection territoriale",
        "screenshot": "commune_Verfeil_resultats.png",
        "description": (
            "Analyse exhaustive d'une commune entière pour identifier tous les sites "
            "potentiels photovoltaïques en une seule requête."
        ),
        "bullets": [
            "Analyse simultanée : parcelles RPG · toitures · parkings · friches industrielles",
            "Résultats streamés en temps réel via SSE (Server-Sent Events)",
            "Filtres cumulables : surface min/max · distance réseau BT/HTA · type de culture",
            "Import Enedis HTA : postes HTB/HTA avec capacités raccordement disponibles",
            "Scoring automatique des sites (0-100) selon critères pondérables",
            "Export tableau Excel/CSV des résultats filtrés",
            "Génération automatique de carte de synthèse (Folium)",
            "Affichage données ZAER (Zones d'Accélération EnR) loi AER 2023",
        ],
        "tech": "SSE · GeoPandas · Shapely · Folium · API IGN · API Enedis · OpenStreetMap",
    },

    # ── SLIDE 8 : CRM Dashboard ──────────────────────────────────────────────
    {
        "type": "feature",
        "category": "CRM – GESTION COMMERCIALE",
        "title": "CRM – Tableau de bord commercial",
        "screenshot": "CRM_Tableau_de_bord.png",
        "description": (
            "CRM intégré nativement dans la plateforme pour gérer l'intégralité "
            "du cycle de vente, de la prospection à la mise en service."
        ),
        "bullets": [
            "Pipeline visuel par étapes : Prospect → Qualifié → Étude → Devis → Signé → Chantier",
            "Fiche prospect complète : coordonnées, adresse, puissance envisagée, statut",
            "Historique des interactions (appels, emails, visites, documents)",
            "Gestion des tâches et relances avec rappels automatiques",
            "Tags et catégorisation multi-critères des prospects",
            "Import en masse (CSV/Excel) de listes de prospects",
            "Intégration directe des résultats de prospection cartographique",
            "API REST complète pour intégrations tierces (webhooks)",
        ],
        "tech": "SQLite/PostgreSQL · Flask-SQLAlchemy · REST API JSON · pandas",
    },

    # ── SLIDE 9 : CRM Stats ──────────────────────────────────────────────────
    {
        "type": "feature",
        "category": "CRM – STATISTIQUES",
        "title": "CRM – Statistiques & KPI",
        "screenshot": "CRM_Statistiques.png",
        "screenshot2": "CRM_Calendrier.png",
        "description": (
            "Tableaux de bord analytiques pour piloter l'activité commerciale "
            "et mesurer la performance de l'équipe."
        ),
        "bullets": [
            "Graphiques interactifs : taux de conversion par étape, CA prévisionnel",
            "KPIs temps réel : prospects actifs, devis en cours, puissance MW signée",
            "Évolution mensuelle du pipeline avec comparaison N-1",
            "Répartition géographique des projets (carte de chaleur)",
            "Calendrier partagé : RDV, visites terrain, jalons chantier",
            "Exports PDF automatisés des rapports de performance",
            "Alertes configurables (prospect inactif, délai relance dépassé)",
            "Tableau de bord personnalisable par utilisateur",
        ],
        "tech": "Chart.js · FullCalendar.js · pandas · reportlab · matplotlib",
    },

    # ── SLIDE 10 : Calpinage PV ──────────────────────────────────────────────
    {
        "type": "feature",
        "category": "CALPINAGE PHOTOVOLTAÏQUE",
        "title": "Outil de calpinage PV 3D sur toiture",
        "screenshot": "calpinage_652.png",
        "description": (
            "Outil de dimensionnement photovoltaïque avancé permettant de dessiner "
            "les panneaux sur la toiture et calculer la production simulée."
        ),
        "bullets": [
            "Dessin interactif des zones de panneaux sur fond satellite haute résolution",
            "Calcul automatique du nombre de modules selon format sélectionné",
            "Simulation production PVGIS avec prise en compte orientation/inclinaison",
            "Calcul perte ombrage (via données DSM LiDAR IGN)",
            "Dimensionnement onduleur automatique selon puissance crête",
            "Calcul autoconsommation + injection réseau avec courbe de charge",
            "Génération du plan de masse PDF exportable (échelle, légende, cartouche)",
            "Conformité RE2020 et normes NF C 15-100",
        ],
        "tech": "Leaflet.js · PVGIS API · LiDAR IGN COPC · Canvas API · jsPDF",
    },

    # ── SLIDE 11 : Autoconsommation ──────────────────────────────────────────
    {
        "type": "feature",
        "category": "SIMULATION ÉNERGÉTIQUE",
        "title": "Simulation autoconsommation & autoconsommation collective",
        "screenshot": "calpinage_652_autoconso.png",
        "screenshot2": "Autoconsommation_Collective.png",
        "description": (
            "Module de simulation complet pour projets d'autoconsommation individuelle "
            "et opérations d'autoconsommation collective (ACC)."
        ),
        "bullets": [
            "Saisie consommation annuelle + profil de consommation (résidentiel/tertiaire/agricole)",
            "Simulation heure par heure sur 8760h (données météo TMY PVGIS)",
            "Calcul taux d'autoconsommation et taux d'autoproduction",
            "Bilan économique sur 25 ans : économies, TRI, VAN, temps de retour",
            "ACC : gestion des participants, clés de répartition personnalisées",
            "Conformité réglementaire loi énergie 2015 et ordonnance ACC 2021",
            "Export rapport PDF complet avec graphiques de production/consommation",
            "Optimisation dimensionnement selon objectif (autoconso max / TRI max)",
        ],
        "tech": "PVGIS API · numpy · pandas · matplotlib · reportlab",
    },

    # ── SLIDE 12 : Rapport commune ───────────────────────────────────────────
    {
        "type": "feature",
        "category": "RAPPORTS TECHNIQUES",
        "title": "Rapport commune – Analyse territoriale complète",
        "screenshot": "rapport_commune_Verfeil_debut.png",
        "description": (
            "Rapport automatisé complet pour une commune : synthèse cartographique, "
            "données énergétiques, urbanistiques et réseaux."
        ),
        "bullets": [
            "Carte de synthèse interactive avec tous les sites identifiés",
            "Tableau récapitulatif : surfaces RPG, toitures, parkings avec potentiel kWc",
            "Données réseau : postes HTA avec puissance disponible (source Enedis)",
            "Contraintes PLU : zones NC, AU, U filtrées avec surfaces constructibles",
            "Données ZAER (Zone d'Accélération EnR) et DUP existantes",
            "Analyse GéoRisques : sismicité, inondations, retrait-gonflement argiles",
            "Données socio-économiques : nombre d'entreprises (SIRENE), exploitants",
            "Export PDF multi-pages avec mise en page professionnelle",
        ],
        "tech": "SSE streaming · GeoPandas · Folium · reportlab · API IGN · GéoRisques",
    },

    # ── SLIDE 13 : Rapport point GPS ─────────────────────────────────────────
    {
        "type": "feature",
        "category": "RAPPORTS TECHNIQUES",
        "title": "Rapport point GPS – Analyse parcellaire détaillée",
        "screenshot": "rapport_point_resultats.png",
        "description": (
            "Rapport de synthèse exhaustif pour un point GPS précis : "
            "toutes les données nécessaires à l'instruction d'un projet PV."
        ),
        "bullets": [
            "Identification cadastrale automatique (section, numéro, surface, propriétaire)",
            "Données PLU : zone, règlement applicable, COS, hauteur max",
            "Potentiel solaire PVGIS : irradiation GHI/DNI, P50/P90, production simulée",
            "Distance et capacité des postes Enedis HTA et BT environnants",
            "Données propriétaires SIREN/SIRET avec activité principale",
            "Risques naturels complets (GéoRisques) : 12 typologies de risques",
            "Servitudes et contraintes (lignes HT, aérodrome, sites classés, Natura 2000)",
            "Génération PDF signable pour instruction administrative",
        ],
        "tech": "API Géoportail · PVGIS · API SIRENE · GéoRisques · Enedis · reportlab",
    },

    # ── SLIDE 14 : Visite technique ──────────────────────────────────────────
    {
        "type": "feature",
        "category": "MODULES SPÉCIALISÉS",
        "title": "Visite technique – Relevés terrain",
        "screenshot": "Visite_Technique.png",
        "screenshot2": "Pre-etude_Proposition.png",
        "description": (
            "Formulaire structuré de visite technique sur site avec génération "
            "automatique du rapport de pré-étude et de la proposition commerciale."
        ),
        "bullets": [
            "Relevés toiture : matériau, pente, orientation, état, surface utile",
            "Installation électrique : puissance souscrite, disjoncteur, compteur (Linky)",
            "Photos terrain intégrées directement depuis mobile (PWA-ready)",
            "Calcul automatique puissance installable + production estimée",
            "Génération proposition commerciale PDF avec prix, financement, économies",
            "Devis automatique selon tarif installateur paramétrable",
            "Envoi email direct au prospect depuis l'interface",
            "Archivage automatique dans la fiche CRM du prospect",
        ],
        "tech": "Flask-WTF · reportlab · python-docx · SMTP · Jinja2 templates",
    },

    # ── SLIDE 15 : Suivi de chantier ─────────────────────────────────────────
    {
        "type": "feature",
        "category": "MODULES SPÉCIALISÉS",
        "title": "Suivi de chantier – Module DOE & conformité",
        "screenshot": "Suivi_de_Chantier.png",
        "description": (
            "Module complet de suivi de l'installation photovoltaïque depuis "
            "la commande jusqu'à la mise en service et la déclaration CONSUEL."
        ),
        "bullets": [
            "Planning de chantier avec jalons : commande → pose → raccordement → consuel → MES",
            "Génération automatique du DOE (Dossier des Ouvrages Exécutés)",
            "Rapport de conformité IEC 62446 (test d'isolement, résistance terre, mesure IV)",
            "Gestion non-conformités (NCF) avec photos et suivi de levée",
            "Schéma unifilaire automatique selon configuration installée",
            "Génération CERFA de déclaration préalable de travaux",
            "Suivi CONSUEL : envoi dossier, date visite, attestation de conformité",
            "Notification automatique Enedis pour mise en service (contrat S21)",
        ],
        "tech": "python-docx · reportlab · PIL · Jinja2 · email · SQLAlchemy",
    },

    # ── SLIDE 16 : LiDAR HD ──────────────────────────────────────────────────
    {
        "type": "feature",
        "category": "ANALYSE AVANCÉE",
        "title": "Analyse LiDAR HD – Modélisation 3D des toitures",
        "screenshot": "Plan_LiDAR_HD.png",
        "description": (
            "Accès direct au nuage de points LiDAR HD de l'IGN pour analyse "
            "fine de la morphologie des toitures et calcul de masques solaires."
        ),
        "bullets": [
            "Streaming direct des tuiles COPC (Cloud-Optimized Point Cloud) de l'IGN",
            "Résolution : 10 pts/m² sur 95% du territoire métropolitain",
            "Extraction automatique des plans de toiture (segmentation planaire)",
            "Calcul précis de la pente et de l'orientation de chaque pan",
            "Modèle Numérique de Surface (MNS) pour calcul d'ombrage",
            "Visualisation 3D interactive dans le navigateur (Three.js / WebGL)",
            "Export DXF / STL pour usage CAO ou intégration BIM",
            "Comparaison avant/après pour suivi post-installation",
        ],
        "tech": "laspy[lazrs] · COPC streaming · HTTP Range Requests · Three.js · NumPy · rasterio",
    },

    # ── SLIDE 17 : AO PV Bâtiment ────────────────────────────────────────────
    {
        "type": "feature",
        "category": "MODULES SPÉCIALISÉS",
        "title": "AO PV Bâtiment – Appels d'offres CRE",
        "screenshot": "AO_PV_Batiment.png",
        "description": (
            "Module dédié à la constitution des dossiers de réponse aux appels d'offres "
            "de la Commission de Régulation de l'Énergie (CRE) pour le bâtiment."
        ),
        "bullets": [
            "Compatibilité AO CRE4 (bâtiments, hangars, ombrières, parkings)",
            "Calcul automatique de la prime à l'énergie selon barème CRE actualisé",
            "Vérification des critères d'éligibilité (surface, puissance, date depermis)",
            "Génération du formulaire Cerfa 12648*01 pré-rempli",
            "Constitution automatique du dossier technique complet",
            "Import des données PVGIS pour attestation de production prévisionnelle",
            "Suivi du statut de l'appel d'offres et des délais réglementaires",
            "Base de données des tarifs CRE avec historique",
        ],
        "tech": "reportlab · pandas · PVGIS · python-docx · openpyxl",
    },

    # ── SLIDE 18 : Helia IA ──────────────────────────────────────────────────
    {
        "type": "feature",
        "category": "INTELLIGENCE ARTIFICIELLE",
        "title": "Helia IA – Assistant solaire intelligent",
        "screenshot": "helia_ia.png",
        "screenshot2": "helia_interface.png",
        "description": (
            "Assistant IA spécialisé dans l'énergie solaire photovoltaïque, "
            "intégré nativement dans la plateforme pour assister les utilisateurs en temps réel."
        ),
        "bullets": [
            "Modèle LLM Groq (LLama 3.1 70B) avec contexte solaire spécialisé",
            "Analyse automatique des données du projet en cours (puissance, surface, localisation)",
            "Réponses aux questions réglementaires (CRE, CONSUEL, urbanisme PV)",
            "Recommandations techniques : onduleur, câblage, protection parafoudre",
            "Calcul de rentabilité expliqué en langage naturel",
            "Génération de compte-rendus de visite technique à la dictée",
            "Mode audit : analyse d'une installation existante sur photo",
            "API REST /api/helia/chat pour intégration dans workflows externes",
        ],
        "tech": "Groq API · LLama 3.1 70B · Streaming SSE · Flask Blueprint",
    },

    # ── SLIDE 19 : Administration ────────────────────────────────────────────
    {
        "type": "feature",
        "category": "ADMINISTRATION",
        "title": "Backoffice & gestion des utilisateurs",
        "screenshot": "Administration.png",
        "screenshot2": "Parametrages_CRM.png",
        "description": (
            "Interface d'administration complète pour gérer les utilisateurs, "
            "les licences et paramétrer l'ensemble de la plateforme."
        ),
        "bullets": [
            "Gestion des utilisateurs : création, modification, suspension, rôles",
            "Licences et abonnements : activation/désactivation, durée, quotas",
            "Monitoring des connexions et activité par utilisateur",
            "Paramétrage CRM : étapes pipeline personnalisées, champs custom",
            "Configuration des tarifs installateur (matériel, pose, BT/HTA)",
            "Gestion des templates de documents (propositions, rapports)",
            "Logs d'audit : toutes les actions tracées avec timestamp/IP",
            "Sauvegarde automatique base de données (pg_dump / SQLite backup)",
        ],
        "tech": "Flask-Admin · SQLAlchemy · bcrypt · python-dotenv · PostgreSQL",
    },

    # ── SLIDE 20 : Architecture technique ───────────────────────────────────
    {
        "type": "tech",
        "title": "Architecture technique",
        "layers": [
            ("Frontend", [
                "Leaflet.js 1.9 – cartographie interactive",
                "Vanilla JS + Bootstrap 5 – UI composants",
                "Three.js / WebGL – rendu 3D LiDAR",
                "Chart.js – graphiques analytics",
                "FullCalendar.js – planning CRM",
                "SSE (EventSource) – streaming temps réel",
            ]),
            ("Backend", [
                "Python 3.11 · Flask 3.1 – serveur applicatif",
                "GeoPandas 0.14 · Shapely 2.0 – calculs géospatiaux",
                "laspy[lazrs] – lecture nuages de points LiDAR COPC",
                "reportlab + python-docx – génération documents",
                "Groq SDK – intégration LLM IA",
                "Stripe SDK – paiement en ligne",
            ]),
            ("Données & APIs", [
                "IGN Géoportail – fond de carte, cadastre, ortho",
                "PVGIS (JRC) – irradiation et simulation PV",
                "API BAN – géocodage adresses françaises",
                "GéoRisques – risques naturels et technologiques",
                "API SIRENE (INSEE) – données entreprises",
                "Enedis Open Data – postes électriques",
            ]),
            ("Infrastructure", [
                "Railway.app – hébergement cloud PaaS (Docker)",
                "PostgreSQL – base de données production",
                "GeoServer – serveur WMS/WFS couches géo",
                "Gunicorn – serveur WSGI production",
                "Redis – cache sessions (optionnel)",
                "GitHub CI/CD – déploiement automatique",
            ]),
        ],
    },

    # ── SLIDE 21 : Statut & abonnements ─────────────────────────────────────
    {
        "type": "twoshot",
        "category": "MONITORING & COMMERCIAL",
        "title": "Monitoring système & Offres d'abonnement",
        "screenshot": "Statut_Systeme.png",
        "screenshot2": "Abonnements.png",
        "desc_left": "Endpoint /health – état en temps réel de tous les services",
        "desc_right": "Offres Starter / Pro / Enterprise avec paiement Stripe",
        "bullets_left": [
            "Statut Flask, GeoServer, base de données",
            "Uptime et latence des APIs tierces",
            "Monitoring mémoire et CPU (Railway metrics)",
        ],
        "bullets_right": [
            "Abonnement mensuel / annuel via Stripe",
            "Essai gratuit 14 jours sans CB",
            "Facturation automatique + portail client Stripe",
        ],
    },

    # ── SLIDE 22 : Conclusion ────────────────────────────────────────────────
    {
        "type": "conclusion",
        "title": "HeliaPV – La plateforme complète du développeur solaire",
        "points": [
            "🗺️  Cartographie multi-couches avec scoring automatique des sites",
            "🔍  Prospection territoriale complète en quelques secondes",
            "📊  CRM intégré du premier contact à la mise en service",
            "🧮  Calpinage 3D + simulation PVGIS + plan de masse PDF",
            "📄  Génération automatique de tous les documents administratifs",
            "🤖  Assistant IA solaire spécialisé (Helia – LLama 3.1 70B)",
            "📡  Données LiDAR HD IGN pour analyse fine des toitures",
            "☁️  SaaS cloud-native, déployable en 1 commande (Railway/Docker)",
        ],
        "url": "https://app.heliapv.fr",
        "contact": "contact@heliapv.fr",
    },
]


# ─────────────────────────────────────────────────────────────────────────────
# FONCTIONS UTILITAIRES
# ─────────────────────────────────────────────────────────────────────────────

def rgb(r, g, b):
    return RGBColor(r, g, b)


def set_bg(slide, color: RGBColor):
    """Remplit le fond de la slide avec une couleur unie."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color, alpha=None):
    """Ajoute un rectangle coloré."""
    shape = slide.shapes.add_shape(
        pptx.util.MSO_SHAPE_TYPE if False else 1,  # MSO_AUTO_SHAPE_TYPE.RECTANGLE = 1
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_text(slide, text, left, top, width, height,
             font_size=18, bold=False, color=C_WHITE,
             align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Segoe UI"
    return txb


def add_image_safe(slide, img_path, left, top, width, height):
    """Ajoute une image si elle existe, sinon un placeholder."""
    path = SCREENSHOTS / img_path if img_path else None
    if path and path.exists() and path.stat().st_size > 10000:
        try:
            return slide.shapes.add_picture(str(path), left, top, width, height)
        except Exception as e:
            print(f"  ⚠️  Image {img_path} impossible à charger: {e}")
    # Placeholder grisé
    ph = add_rect(slide, left, top, width, height, rgb(0x1C, 0x23, 0x33))
    add_text(slide, f"📷  {img_path or 'N/A'}\n(capture non disponible)",
             left + Inches(0.2), top + height//2 - Pt(20),
             width - Inches(0.4), Inches(0.8),
             font_size=10, color=C_MUTED, align=PP_ALIGN.CENTER)
    return ph


def add_bullet_list(slide, bullets, left, top, width, height,
                    font_size=11, color=C_WHITE, accent=C_ACCENT):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf  = txb.text_frame
    tf.word_wrap = True
    first = True
    for b in bullets:
        if first:
            para = tf.paragraphs[0]
            first = False
        else:
            para = tf.add_paragraph()
        para.space_before = Pt(2)
        # Puce colorée
        r1 = para.add_run()
        r1.text = "▸  "
        r1.font.color.rgb = accent
        r1.font.size = Pt(font_size)
        r1.font.name = "Segoe UI"
        # Texte
        r2 = para.add_run()
        r2.text = b
        r2.font.color.rgb = color
        r2.font.size = Pt(font_size)
        r2.font.name = "Segoe UI"
    return txb


def add_category_pill(slide, text, left, top):
    """Badge catégorie accent."""
    w = Inches(3.2)
    h = Inches(0.28)
    bg = add_rect(slide, left, top, w, h, C_ACCENT)
    add_text(slide, text, left + Inches(0.1), top, w - Inches(0.1), h,
             font_size=8, bold=True, color=C_BG_DARK)


def add_divider(slide, top, color=C_ACCENT):
    ln = slide.shapes.add_shape(1, Inches(0.4), top, Inches(0.05), Inches(0.03))
    ln.fill.solid()
    ln.fill.fore_color.rgb = color
    ln.line.fill.background()
    # Ligne horizontale
    ln2 = slide.shapes.add_shape(1, Inches(0.4), top, Inches(12.5), Inches(0.03))
    ln2.fill.solid()
    ln2.fill.fore_color.rgb = color
    ln2.line.fill.background()


def add_tech_badge(slide, text, left, top):
    w = max(Inches(1.6), Pt(len(text) * 6))
    h = Inches(0.22)
    bg = add_rect(slide, left, top, w, h, rgb(0x1C, 0x23, 0x33))
    add_text(slide, text, left + Inches(0.05), top, w - Inches(0.05), h,
             font_size=7.5, color=C_ACCENT2)
    return w


# ─────────────────────────────────────────────────────────────────────────────
# GÉNÉRATEURS DE SLIDES
# ─────────────────────────────────────────────────────────────────────────────

def make_cover(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_bg(slide, C_BG_DARK)

    # Bande accent gauche
    add_rect(slide, 0, 0, Inches(0.1), H, C_ACCENT)

    # Soleil décoratif (cercle)
    sun = slide.shapes.add_shape(1, Inches(9.5), Inches(0.5), Inches(3.5), Inches(3.5))
    sun.fill.solid()
    sun.fill.fore_color.rgb = rgb(0xF5, 0x9E, 0x0B)
    sun.line.fill.background()
    # translucide simulé par couleur plus sombre
    sun2 = slide.shapes.add_shape(1, Inches(9.6), Inches(0.6), Inches(3.3), Inches(3.3))
    sun2.fill.solid()
    sun2.fill.fore_color.rgb = C_BG_DARK
    sun2.line.fill.background()

    # Logo texte
    add_text(slide, "☀️  HeliaPV", Inches(0.5), Inches(1.2), Inches(5), Inches(0.8),
             font_size=14, bold=True, color=C_ACCENT)

    # Titre principal
    txb = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(8.5), Inches(1.8))
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = data["title"]
    r.font.size = Pt(40)
    r.font.bold = True
    r.font.color.rgb = C_WHITE
    r.font.name = "Segoe UI"

    # Sous-titre
    add_text(slide, data["subtitle"], Inches(0.5), Inches(3.8), Inches(9), Inches(0.6),
             font_size=16, color=C_MUTED)

    # Tagline
    add_text(slide, data["tagline"], Inches(0.5), Inches(4.5), Inches(9), Inches(0.4),
             font_size=12, color=C_ACCENT2, italic=True)

    # Version
    add_text(slide, data["version"], Inches(0.5), Inches(6.8), Inches(4), Inches(0.3),
             font_size=10, color=C_MUTED)

    # URL
    add_text(slide, "app.heliapv.fr", Inches(9.5), Inches(6.8), Inches(3.5), Inches(0.3),
             font_size=10, color=C_ACCENT, align=PP_ALIGN.RIGHT)

    # Bande basse
    add_rect(slide, 0, H - Inches(0.08), W, Inches(0.08), C_ACCENT)

    print("  ✅  Slide 1 : Couverture")


def make_toc(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, C_BG_DARK)
    add_rect(slide, 0, 0, Inches(0.08), H, C_ACCENT)

    add_text(slide, data["title"], Inches(0.4), Inches(0.3), Inches(12), Inches(0.6),
             font_size=28, bold=True, color=C_WHITE)
    add_rect(slide, Inches(0.4), Inches(0.95), Inches(12.5), Inches(0.04), C_ACCENT)

    cols = [data["sections"][:5], data["sections"][5:]]
    for ci, col in enumerate(cols):
        x = Inches(0.5 + ci * 6.4)
        for i, (num, title, sub) in enumerate(col):
            y = Inches(1.15 + i * 1.1)
            # Numéro
            add_rect(slide, x, y, Inches(0.45), Inches(0.45), C_ACCENT)
            add_text(slide, num, x, y + Inches(0.05), Inches(0.45), Inches(0.35),
                     font_size=14, bold=True, color=C_BG_DARK, align=PP_ALIGN.CENTER)
            # Titre
            add_text(slide, title, x + Inches(0.55), y, Inches(5.7), Inches(0.3),
                     font_size=12, bold=True, color=C_WHITE)
            # Sous-titre
            add_text(slide, sub, x + Inches(0.55), y + Inches(0.3), Inches(5.7), Inches(0.25),
                     font_size=9, color=C_MUTED, italic=True)

    print("  ✅  Slide 2 : Sommaire")


def make_feature(prs, data, slide_n):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, C_BG_DARK)
    add_rect(slide, 0, 0, Inches(0.08), H, C_ACCENT)
    add_rect(slide, 0, H - Inches(0.08), W, Inches(0.08), C_ACCENT2)

    # Catégorie
    add_category_pill(slide, data.get("category", ""), Inches(0.3), Inches(0.2))

    # Titre
    add_text(slide, data["title"], Inches(0.3), Inches(0.55), Inches(12.7), Inches(0.65),
             font_size=22, bold=True, color=C_WHITE)
    add_rect(slide, Inches(0.3), Inches(1.2), Inches(12.7), Inches(0.03), C_ACCENT)

    # Description
    add_text(slide, data.get("description", ""), Inches(0.3), Inches(1.28),
             Inches(12.7), Inches(0.5),
             font_size=10, color=C_MUTED, italic=True)

    has_two_shots = "screenshot2" in data and data["screenshot2"]

    if has_two_shots:
        # Deux images côte à côte
        add_image_safe(slide, data["screenshot"],  Inches(0.3),  Inches(1.85), Inches(5.5), Inches(3.5))
        add_image_safe(slide, data["screenshot2"], Inches(5.95), Inches(1.85), Inches(5.5), Inches(3.5))
        # Bullets en bas
        add_bullet_list(slide, data.get("bullets", []), Inches(0.3), Inches(5.5),
                        Inches(12.7), Inches(1.7), font_size=9.5)
    else:
        # Une image à gauche, bullets à droite
        add_image_safe(slide, data.get("screenshot"), Inches(0.3), Inches(1.85),
                       Inches(7.2), Inches(4.8))
        # Bullets
        add_bullet_list(slide, data.get("bullets", []), Inches(7.7), Inches(1.85),
                        Inches(5.3), Inches(4.0), font_size=10.5)
        # Badge tech
        if data.get("tech"):
            add_rect(slide, Inches(7.7), Inches(6.0), Inches(5.3), Inches(0.25),
                     rgb(0x1C, 0x23, 0x33))
            add_text(slide, f"🛠  {data['tech']}", Inches(7.8), Inches(6.0),
                     Inches(5.1), Inches(0.25), font_size=7.5, color=C_ACCENT2)

    # Numéro de slide
    add_text(slide, str(slide_n), Inches(12.9), Inches(7.1), Inches(0.4), Inches(0.3),
             font_size=10, color=C_MUTED, align=PP_ALIGN.RIGHT)

    print(f"  ✅  Slide {slide_n} : {data['title'][:50]}")


def make_twoshot(prs, data, slide_n):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, C_BG_DARK)
    add_rect(slide, 0, 0, Inches(0.08), H, C_ACCENT)

    add_category_pill(slide, data.get("category", ""), Inches(0.3), Inches(0.2))
    add_text(slide, data["title"], Inches(0.3), Inches(0.55), Inches(12.7), Inches(0.6),
             font_size=22, bold=True, color=C_WHITE)
    add_rect(slide, Inches(0.3), Inches(1.15), Inches(12.7), Inches(0.03), C_ACCENT)

    # Image gauche
    add_image_safe(slide, data["screenshot"],  Inches(0.3),  Inches(1.25), Inches(6.0), Inches(4.2))
    add_text(slide, data.get("desc_left", ""), Inches(0.3), Inches(5.5), Inches(6.0), Inches(0.3),
             font_size=9, color=C_MUTED, italic=True, align=PP_ALIGN.CENTER)

    # Image droite
    add_image_safe(slide, data["screenshot2"], Inches(6.5),  Inches(1.25), Inches(6.5), Inches(4.2))
    add_text(slide, data.get("desc_right", ""), Inches(6.5), Inches(5.5), Inches(6.5), Inches(0.3),
             font_size=9, color=C_MUTED, italic=True, align=PP_ALIGN.CENTER)

    # Bullets gauche + droite
    add_bullet_list(slide, data.get("bullets_left", []),  Inches(0.3),  Inches(5.9),
                    Inches(5.8), Inches(1.3), font_size=9)
    add_bullet_list(slide, data.get("bullets_right", []), Inches(6.5), Inches(5.9),
                    Inches(6.5), Inches(1.3), font_size=9)

    add_text(slide, str(slide_n), Inches(12.9), Inches(7.1), Inches(0.4), Inches(0.3),
             font_size=10, color=C_MUTED, align=PP_ALIGN.RIGHT)
    print(f"  ✅  Slide {slide_n} : {data['title'][:50]}")


def make_tech(prs, data, slide_n):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, C_BG_DARK)
    add_rect(slide, 0, 0, Inches(0.08), H, C_ACCENT)

    add_text(slide, data["title"], Inches(0.3), Inches(0.2), Inches(12.7), Inches(0.6),
             font_size=26, bold=True, color=C_WHITE)
    add_rect(slide, Inches(0.3), Inches(0.82), Inches(12.7), Inches(0.03), C_ACCENT)

    layers = data["layers"]
    cols = 4
    col_w = Inches(3.1)
    for i, (layer_name, items) in enumerate(layers):
        x = Inches(0.3 + i * 3.2)
        # En-tête colonne
        add_rect(slide, x, Inches(0.9), col_w - Inches(0.1), Inches(0.35), C_ACCENT)
        add_text(slide, layer_name, x + Inches(0.05), Inches(0.9),
                 col_w - Inches(0.15), Inches(0.35),
                 font_size=11, bold=True, color=C_BG_DARK)
        # Items
        for j, item in enumerate(items):
            y = Inches(1.35 + j * 0.95)
            add_rect(slide, x, y, col_w - Inches(0.1), Inches(0.85), rgb(0x14, 0x1C, 0x2E))
            # Icône
            add_text(slide, "◆", x + Inches(0.05), y + Inches(0.08),
                     Inches(0.2), Inches(0.3), font_size=7, color=C_ACCENT)
            add_text(slide, item, x + Inches(0.22), y + Inches(0.05),
                     col_w - Inches(0.35), Inches(0.7),
                     font_size=9, color=C_WHITE)

    add_text(slide, str(slide_n), Inches(12.9), Inches(7.1), Inches(0.4), Inches(0.3),
             font_size=10, color=C_MUTED, align=PP_ALIGN.RIGHT)
    print(f"  ✅  Slide {slide_n} : {data['title']}")


def make_conclusion(prs, data, slide_n):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, C_BG_DARK)

    # Bandes décoratives
    add_rect(slide, 0, 0, Inches(0.08), H, C_ACCENT)
    add_rect(slide, 0, H - Inches(0.08), W, Inches(0.08), C_ACCENT)
    add_rect(slide, 0, 0, W, Inches(0.08), C_ACCENT2)

    add_text(slide, data["title"], Inches(0.4), Inches(0.3), Inches(12.5), Inches(0.8),
             font_size=24, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_rect(slide, Inches(2.5), Inches(1.15), Inches(8.3), Inches(0.04), C_ACCENT)

    # Points forts en grille 2x4
    pts = data["points"]
    for i, pt in enumerate(pts):
        col = i % 2
        row = i // 2
        x = Inches(0.5 + col * 6.4)
        y = Inches(1.35 + row * 1.3)
        # Carte
        add_rect(slide, x, y, Inches(6.0), Inches(1.1), rgb(0x14, 0x1C, 0x2E))
        add_text(slide, pt, x + Inches(0.2), y + Inches(0.15), Inches(5.6), Inches(0.8),
                 font_size=12, color=C_WHITE)

    # URL et contact
    add_text(slide, f"🌐  {data['url']}", Inches(0.4), Inches(6.8), Inches(6), Inches(0.3),
             font_size=13, bold=True, color=C_ACCENT)
    add_text(slide, f"✉️  {data['contact']}", Inches(7.0), Inches(6.8), Inches(6), Inches(0.3),
             font_size=13, color=C_MUTED, align=PP_ALIGN.RIGHT)

    print(f"  ✅  Slide {slide_n} : Conclusion")


# ─────────────────────────────────────────────────────────────────────────────
# MAIN
# ─────────────────────────────────────────────────────────────────────────────

def generate():
    print("\n" + "═"*60)
    print("  ☀️   Génération du PowerPoint HeliaPV / AgriWeb")
    print("═"*60 + "\n")

    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    slide_n = 1
    for data in SLIDES:
        t = data["type"]
        if t == "cover":
            make_cover(prs, data)
        elif t == "toc":
            make_toc(prs, data)
        elif t == "feature":
            make_feature(prs, data, slide_n)
        elif t == "twoshot":
            make_twoshot(prs, data, slide_n)
        elif t == "tech":
            make_tech(prs, data, slide_n)
        elif t == "conclusion":
            make_conclusion(prs, data, slide_n)
        slide_n += 1

    prs.save(str(OUTPUT))
    print(f"\n  ✅  Fichier sauvegardé : {OUTPUT.absolute()}")
    print(f"  📊  {slide_n - 1} slides générées\n")


if __name__ == "__main__":
    generate()
